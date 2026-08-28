from __future__ import annotations

import re
import sys
import zipfile
from pathlib import Path

SCRIPT_DIR = Path(__file__).resolve().parent
PRACTICE_DIR = SCRIPT_DIR.parent
WORKSPACE = Path(__file__).resolve().parents[4]
sys.path.insert(0, str(WORKSPACE / "pptdeps"))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


SLIDE_W = 13.333
SLIDE_H = 7.5
FONT = "Malgun Gothic"


def rgb(value: str) -> RGBColor:
    value = value.lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def new_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    return prs


def blank(prs: Presentation, background: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(background)
    return slide


def box(slide, x, y, w, h, fill, radius=True, line=None, line_width=1):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    if line:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def circle(slide, x, y, d, fill, line=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    if line:
        shape.line.color.rgb = rgb(line)
    else:
        shape.line.fill.background()
    return shape


def line(slide, x1, y1, x2, y2, color, width=2):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x1),
        Inches(y1),
        Inches(max(x2 - x1, 0.01)),
        Inches(max(y2 - y1, 0.01)),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.line.fill.background()
    if y2 == y1:
        shape.height = Pt(width)
    if x2 == x1:
        shape.width = Pt(width)
    return shape


def text(
    slide,
    value,
    x,
    y,
    w,
    h,
    size=20,
    color="#111827",
    bold=False,
    align="left",
    valign="middle",
    margin=0.05,
):
    shape = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = {
        "top": MSO_ANCHOR.TOP,
        "middle": MSO_ANCHOR.MIDDLE,
        "bottom": MSO_ANCHOR.BOTTOM,
    }[valign]
    paragraph = frame.paragraphs[0]
    paragraph.alignment = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }[align]
    run = paragraph.add_run()
    run.text = str(value)
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return shape


def pill(slide, label, x, y, w, fill, color, size=11):
    box(slide, x, y, w, 0.34, fill, True)
    text(slide, label, x, y, w, 0.34, size, color, True, "center")


def header(slide, eyebrow, title_value, page, palette, dark=False):
    muted = palette["muted_dark"] if dark else palette["muted"]
    primary = palette["text_dark"] if dark else palette["text"]
    pill(slide, eyebrow, 0.65, 0.40, 1.65, palette["accent"], palette["accent_text"], 10)
    text(slide, title_value, 0.65, 0.88, 11.5, 0.74, 26, primary, True, "left")
    text(slide, f"{page:02d}", 12.05, 0.42, 0.55, 0.34, 10, muted, True, "right")


def footer(slide, label, palette, dark=False, note=None):
    muted = palette["muted_dark"] if dark else palette["muted"]
    line(slide, 0.65, 7.12, 12.68, 7.12, palette["line_dark"] if dark else palette["line"], 1)
    text(slide, label, 0.65, 7.16, 4.8, 0.20, 8.5, muted)
    if note:
        text(slide, note, 7.0, 7.16, 5.65, 0.20, 8.5, muted, False, "right")


def number_card(slide, x, y, w, h, number, title_value, body, palette, dark=False):
    fill = palette["panel_dark"] if dark else palette["panel"]
    text_color = palette["text_dark"] if dark else palette["text"]
    muted = palette["muted_dark"] if dark else palette["muted"]
    box(slide, x, y, w, h, fill, True, palette["line_dark"] if dark else palette["line"])
    circle(slide, x + 0.25, y + 0.24, 0.48, palette["accent"])
    text(slide, number, x + 0.25, y + 0.24, 0.48, 0.48, 13, palette["accent_text"], True, "center")
    text(slide, title_value, x + 0.25, y + 0.88, w - 0.5, 0.42, 18, text_color, True)
    text(slide, body, x + 0.25, y + 1.34, w - 0.5, h - 1.55, 12.5, muted, False, "left", "top")


BASIC = {
    "bg": "#F7FAFF",
    "panel": "#FFFFFF",
    "panel_dark": "#0F2454",
    "text": "#14213D",
    "text_dark": "#FFFFFF",
    "muted": "#5B6B88",
    "muted_dark": "#BDD0F6",
    "accent": "#2F6BFF",
    "accent2": "#12B8A6",
    "accent_text": "#FFFFFF",
    "line": "#D9E4F6",
    "line_dark": "#29457B",
}

MID = {
    "bg": "#F3F7F8",
    "panel": "#FFFFFF",
    "panel_dark": "#102A38",
    "text": "#102A38",
    "text_dark": "#F5FBFC",
    "muted": "#5F7380",
    "muted_dark": "#A7C0C8",
    "accent": "#13B8A6",
    "accent2": "#3A83F6",
    "accent_text": "#062D2B",
    "line": "#D5E2E5",
    "line_dark": "#28505D",
}

ADV = {
    "bg": "#12151C",
    "panel": "#1A1F2A",
    "panel_dark": "#1A1F2A",
    "text": "#F4F7FB",
    "text_dark": "#F4F7FB",
    "muted": "#9CA7B8",
    "muted_dark": "#9CA7B8",
    "accent": "#42D6E8",
    "accent2": "#9B7BFF",
    "accent_text": "#071316",
    "line": "#303847",
    "line_dark": "#303847",
}


def build_basic():
    prs = new_deck()

    slide = blank(prs, BASIC["bg"])
    box(slide, 8.5, 0, 4.83, 7.5, "#E5EEFF", False)
    circle(slide, 9.05, 1.05, 3.35, "#2F6BFF")
    circle(slide, 10.05, 2.05, 1.35, "#12B8A6")
    line(slide, 9.9, 2.55, 11.4, 2.55, "#FFFFFF", 6)
    line(slide, 10.65, 1.8, 10.65, 3.3, "#FFFFFF", 6)
    pill(slide, "DAY9 · BASIC", 0.75, 0.72, 1.7, BASIC["accent"], "#FFFFFF", 10)
    text(slide, "AI 에이전트\n업무 입문", 0.75, 1.55, 7.0, 1.65, 39, BASIC["text"], True, "left", "top")
    text(slide, "반복 업무를 줄이는 디지털 동료", 0.78, 3.42, 6.2, 0.48, 20, BASIC["muted"])
    text(slide, "기초 실습 결과 예시 · 6 SLIDES", 0.78, 6.63, 4.8, 0.35, 10, BASIC["muted"], True)

    slide = blank(prs, BASIC["bg"])
    header(slide, "01 · CONCEPT", "AI 에이전트는 ‘목표를 끝까지 수행하는 시스템’입니다", 2, BASIC)
    parts = [
        ("목표", "무엇을 끝낼지\n명확히 정의"),
        ("판단", "순서를 계획하고\n다음 행동 선택"),
        ("도구", "파일·검색·스크립트로\n실제 작업 실행"),
    ]
    for i, (title_value, body) in enumerate(parts):
        number_card(slide, 0.75 + i * 4.15, 2.05, 3.72, 2.55, f"0{i+1}", title_value, body, BASIC)
    box(slide, 2.15, 5.25, 9.05, 0.92, "#E7EEFF", True)
    text(slide, "사람의 요청  +  스스로 계획  +  도구 실행  +  결과 검증", 2.15, 5.25, 9.05, 0.92, 20, BASIC["accent"], True, "center")
    footer(slide, "AI 에이전트 업무 입문", BASIC)

    slide = blank(prs, BASIC["bg"])
    header(slide, "02 · PROBLEM", "반복 업무에서는 세 곳에서 시간이 새어 나갑니다", 3, BASIC)
    items = [
        ("01", "찾기", "흩어진 자료와\n최신 버전을 찾는 시간", "매번 다시 검색"),
        ("02", "옮기기", "메일·문서·표 사이\n수작업 복사와 정리", "형식 오류 발생"),
        ("03", "확인하기", "누락·오탈자·규칙을\n사람이 반복 점검", "검토 대기 증가"),
    ]
    for i, (num, title_value, body, tag) in enumerate(items):
        x = 0.75 + i * 4.15
        box(slide, x, 1.92, 3.72, 3.75, "#FFFFFF", True, BASIC["line"])
        text(slide, num, x + 0.28, 2.2, 0.8, 0.48, 25, BASIC["accent"], True)
        text(slide, title_value, x + 0.28, 2.84, 3.0, 0.5, 21, BASIC["text"], True)
        text(slide, body, x + 0.28, 3.55, 3.05, 1.05, 14, BASIC["muted"], False, "left", "top")
        pill(slide, tag, x + 0.28, 5.0, 1.65, "#EEF3FF", BASIC["accent"], 9.5)
    footer(slide, "AI 에이전트 업무 입문", BASIC, note="교육용 상황 예시")

    slide = blank(prs, BASIC["bg"])
    header(slide, "03 · FLOW", "좋은 결과는 ‘요청 → 계획 → 실행 → 검증’에서 나옵니다", 4, BASIC)
    stages = [
        ("1", "요청", "목적·대상·완료 기준"),
        ("2", "계획", "작업 순서와 필요한 도구"),
        ("3", "실행", "문서·데이터·파일 처리"),
        ("4", "검증", "누락·형식·품질 점검"),
    ]
    for i, (num, title_value, body) in enumerate(stages):
        x = 0.72 + i * 3.12
        box(slide, x, 2.2, 2.56, 2.6, "#FFFFFF", True, BASIC["line"])
        circle(slide, x + 0.88, 1.85, 0.8, BASIC["accent"] if i < 3 else BASIC["accent2"])
        text(slide, num, x + 0.88, 1.85, 0.8, 0.8, 18, "#FFFFFF", True, "center")
        text(slide, title_value, x + 0.2, 2.75, 2.16, 0.44, 20, BASIC["text"], True, "center")
        text(slide, body, x + 0.25, 3.4, 2.06, 0.72, 12.5, BASIC["muted"], False, "center", "top")
        if i < 3:
            text(slide, "→", x + 2.58, 3.0, 0.55, 0.6, 23, BASIC["accent"], True, "center")
    box(slide, 1.6, 5.35, 10.1, 0.68, "#E8FBF8", True)
    text(slide, "핵심: ‘무엇을 만들까?’보다 ‘어떻게 완료를 확인할까?’를 먼저 적습니다.", 1.6, 5.35, 10.1, 0.68, 15, "#087C72", True, "center")
    footer(slide, "AI 에이전트 업무 입문", BASIC)

    slide = blank(prs, BASIC["bg"])
    header(slide, "04 · USE CASES", "작고 반복적인 업무부터 맡기면 안전하게 시작할 수 있습니다", 5, BASIC)
    cases = [
        ("A", "문서 요약", "긴 회의록을\n결정·할 일로 정리"),
        ("B", "이메일 초안", "대상과 목적에 맞춘\n1차 문안 작성"),
        ("C", "자료 조사", "질문별 근거와\n출처 후보 수집"),
        ("D", "보고서 정리", "표·문장·목차의\n형식을 일관되게"),
    ]
    for i, (letter, title_value, body) in enumerate(cases):
        x = 0.75 + (i % 2) * 6.18
        y = 1.92 + (i // 2) * 2.18
        box(slide, x, y, 5.65, 1.76, "#FFFFFF", True, BASIC["line"])
        box(slide, x + 0.22, y + 0.28, 0.82, 1.2, BASIC["accent"] if i % 2 == 0 else BASIC["accent2"], True)
        text(slide, letter, x + 0.22, y + 0.28, 0.82, 1.2, 22, "#FFFFFF", True, "center")
        text(slide, title_value, x + 1.28, y + 0.3, 2.6, 0.42, 17, BASIC["text"], True)
        text(slide, body, x + 1.28, y + 0.8, 3.85, 0.7, 12, BASIC["muted"], False, "left", "top")
    footer(slide, "AI 에이전트 업무 입문", BASIC)

    slide = blank(prs, BASIC["bg"])
    header(slide, "05 · START", "7일 안에 ‘작은 성공’을 만드는 실행 계획", 6, BASIC)
    days = [
        ("DAY 1", "업무 1개 선택", "주 2회 이상 반복"),
        ("DAY 2", "완료 기준 작성", "예시 결과 포함"),
        ("DAY 3–4", "초안 실행", "사람이 즉시 검토"),
        ("DAY 5", "실패 조건 보완", "금지·예외 추가"),
        ("DAY 6–7", "재실행·공유", "전후 차이 기록"),
    ]
    line(slide, 1.1, 3.02, 12.18, 3.02, BASIC["line"], 4)
    for i, (day, title_value, body) in enumerate(days):
        x = 0.72 + i * 2.48
        circle(slide, x + 0.72, 2.68, 0.68, BASIC["accent"] if i < 4 else BASIC["accent2"])
        text(slide, str(i + 1), x + 0.72, 2.68, 0.68, 0.68, 14, "#FFFFFF", True, "center")
        text(slide, day, x, 3.52, 2.12, 0.3, 10, BASIC["accent"], True, "center")
        text(slide, title_value, x, 3.9, 2.12, 0.45, 14, BASIC["text"], True, "center")
        text(slide, body, x, 4.42, 2.12, 0.55, 10.5, BASIC["muted"], False, "center", "top")
    box(slide, 2.0, 5.55, 9.35, 0.68, "#14213D", True)
    text(slide, "실습 과제  ·  오늘 반복한 업무 하나를 4단계 프롬프트로 바꾸기", 2.0, 5.55, 9.35, 0.68, 16, "#FFFFFF", True, "center")
    footer(slide, "AI 에이전트 업무 입문", BASIC)

    path = PRACTICE_DIR / "PPT기초.pptx"
    prs.save(path)
    return path


def build_mid():
    prs = new_deck()

    slide = blank(prs, MID["panel_dark"])
    box(slide, 8.1, 0, 5.23, 7.5, "#0B202A", False)
    for i, width in enumerate([3.7, 3.15, 2.55, 1.95]):
        box(slide, 8.85, 1.3 + i * 1.05, width, 0.62, MID["accent"] if i < 2 else MID["accent2"], True)
    pill(slide, "8-WEEK PILOT", 0.78, 0.72, 1.85, MID["accent"], MID["accent_text"], 10)
    text(slide, "AI 에이전트\n업무 자동화 파일럿", 0.78, 1.62, 7.0, 1.75, 36, MID["text_dark"], True, "left", "top")
    text(slide, "팀장·실무 리더 의사결정용 제안", 0.82, 3.62, 5.8, 0.45, 18, MID["muted_dark"])
    text(slide, "내부 검토용 · 가정 기반 예시", 0.82, 6.68, 4.2, 0.3, 10, MID["muted_dark"], True)

    slide = blank(prs, MID["bg"])
    header(slide, "01 · SUMMARY", "8주 동안 3개 프로세스의 ‘초안–검토’ 구간을 검증합니다", 2, MID)
    metrics = [("420h", "월 반복업무"), ("20명", "파일럿 사용자"), ("3개", "대상 프로세스")]
    for i, (value, label) in enumerate(metrics):
        x = 0.75 + i * 2.55
        box(slide, x, 2.02, 2.25, 1.55, "#FFFFFF", True, MID["line"])
        text(slide, value, x + 0.18, 2.28, 1.9, 0.58, 27, MID["text"], True)
        text(slide, label, x + 0.18, 2.92, 1.9, 0.3, 11, MID["muted"])
    box(slide, 8.55, 2.02, 4.02, 3.55, MID["panel_dark"], True)
    text(slide, "오늘 필요한 결정", 8.9, 2.35, 3.25, 0.42, 17, MID["text_dark"], True)
    decisions = ["대상 업무 3개 승인", "실무 사용자 20명 배정", "데이터 접근 범위 확정"]
    for i, item in enumerate(decisions):
        circle(slide, 8.9, 3.15 + i * 0.67, 0.25, MID["accent"])
        text(slide, item, 9.32, 3.04 + i * 0.67, 2.72, 0.45, 12.5, MID["text_dark"], True)
    box(slide, 0.75, 4.05, 7.25, 1.52, "#DFF7F3", True)
    text(slide, "가설", 1.02, 4.3, 0.68, 0.36, 11, "#087C72", True)
    text(slide, "초안 자동화 + 사람 승인으로 속도와 품질을 동시에 개선할 수 있다.", 1.02, 4.72, 6.48, 0.5, 16, MID["text"], True)
    footer(slide, "AI 에이전트 업무 자동화 파일럿", MID, note="가정 기반 예시")

    slide = blank(prs, MID["bg"])
    header(slide, "02 · BOTTLENECK", "수작업·재작업·대기시간이 같은 지점에 겹칩니다", 3, MID)
    stats = [
        ("420", "시간/월", "수작업 취합"),
        ("18", "%", "재작업률"),
        ("2", "일", "검토 SLA"),
    ]
    for i, (value, unit, label) in enumerate(stats):
        x = 0.75 + i * 4.15
        box(slide, x, 1.95, 3.72, 1.52, "#FFFFFF", True, MID["line"])
        text(slide, value, x + 0.25, 2.2, 1.55, 0.58, 27, MID["text"], True)
        text(slide, unit, x + 1.7, 2.34, 0.8, 0.34, 11, MID["muted"], True)
        text(slide, label, x + 0.25, 2.88, 2.8, 0.3, 11, MID["muted"])
    rows = [
        ("자료 찾기", 0.78, "#3A83F6"),
        ("초안 작성", 0.62, "#13B8A6"),
        ("형식 맞추기", 0.48, "#7E93A0"),
        ("검토 대기", 0.86, "#F2A65A"),
    ]
    for i, (label, score, fill) in enumerate(rows):
        y = 4.05 + i * 0.48
        text(slide, label, 0.9, y, 1.25, 0.28, 10.5, MID["muted"], True)
        box(slide, 2.25, y + 0.04, 8.8, 0.22, "#DDE8EA", True)
        box(slide, 2.25, y + 0.04, 8.8 * score, 0.22, fill, True)
        text(slide, f"{int(score * 100)}", 11.2, y, 0.48, 0.28, 10, MID["text"], True, "right")
    footer(slide, "AI 에이전트 업무 자동화 파일럿", MID, note="교육용 가정치")

    slide = blank(prs, MID["bg"])
    header(slide, "03 · ARCHITECTURE", "요청을 받는 순간부터 사람 승인까지 하나의 흐름으로 묶습니다", 4, MID)
    nodes = [
        ("INTAKE", "요청·파일"),
        ("ORCHESTRATOR", "계획·분기"),
        ("SKILLS", "업무 절차"),
        ("SYSTEMS", "문서·데이터"),
        ("REVIEW", "사람 승인"),
    ]
    for i, (label, sub) in enumerate(nodes):
        x = 0.55 + i * 2.55
        fill = MID["panel_dark"] if i in (1, 4) else "#FFFFFF"
        fg = MID["text_dark"] if i in (1, 4) else MID["text"]
        box(slide, x, 2.55, 2.05, 1.65, fill, True, MID["line"])
        text(slide, f"0{i+1}", x + 0.18, 2.75, 0.55, 0.32, 10, MID["accent"], True)
        text(slide, label, x + 0.18, 3.13, 1.7, 0.35, 13, fg, True)
        text(slide, sub, x + 0.18, 3.55, 1.7, 0.28, 10.5, MID["muted_dark"] if i in (1, 4) else MID["muted"])
        if i < 4:
            text(slide, "→", x + 2.05, 3.08, 0.5, 0.5, 20, MID["accent"], True, "center")
    box(slide, 2.1, 4.88, 9.12, 0.82, "#DFF7F3", True)
    text(slide, "자동 실행 범위는 제한하고, 외부 발송·정책 판단은 REVIEW 단계에서 승인합니다.", 2.1, 4.88, 9.12, 0.82, 14, "#087C72", True, "center")
    footer(slide, "AI 에이전트 업무 자동화 파일럿", MID)

    slide = blank(prs, MID["bg"])
    header(slide, "04 · PRIORITY", "효과가 크고 구현이 쉬운 두 업무부터 시작합니다", 5, MID)
    box(slide, 1.2, 1.92, 10.4, 4.48, "#FFFFFF", False, MID["line"])
    line(slide, 6.4, 1.92, 6.4, 6.4, MID["line"], 2)
    line(slide, 1.2, 4.16, 11.6, 4.16, MID["line"], 2)
    text(slide, "효과 높음", 0.55, 2.02, 0.58, 0.4, 9.5, MID["muted"], True, "center")
    text(slide, "효과 낮음", 0.55, 5.78, 0.58, 0.4, 9.5, MID["muted"], True, "center")
    text(slide, "구현 쉬움", 1.2, 6.48, 1.5, 0.3, 9.5, MID["muted"], True)
    text(slide, "구현 어려움", 10.0, 6.48, 1.6, 0.3, 9.5, MID["muted"], True, "right")
    bubbles = [
        (3.2, 2.5, 1.1, "주간 보고", MID["accent"]),
        (4.7, 3.3, 1.0, "VOC 분류", MID["accent"]),
        (7.3, 2.55, 1.2, "제안서 초안", MID["accent2"]),
        (8.9, 4.75, 1.1, "캠페인 QA", "#F2A65A"),
        (5.1, 5.12, 0.95, "지식 검색", "#7E93A0"),
    ]
    for x, y, d, label, fill in bubbles:
        circle(slide, x, y, d, fill)
        text(slide, label, x - 0.15, y + d / 2 - 0.18, d + 0.3, 0.36, 9.5, "#FFFFFF", True, "center")
    pill(slide, "1차 후보", 1.48, 2.15, 1.18, "#DFF7F3", "#087C72", 9)
    footer(slide, "AI 에이전트 업무 자동화 파일럿", MID, note="효과·난이도는 인터뷰 기반 가정")

    slide = blank(prs, MID["bg"])
    header(slide, "05 · OPERATING MODEL", "AI는 초안을 만들고, 책임자는 기준에 따라 승인합니다", 6, MID)
    lanes = [
        ("REQUESTER", "업무 요청", "목표·자료·기한 입력"),
        ("AGENT", "초안 생성", "규칙 적용·근거 기록"),
        ("REVIEWER", "품질 승인", "수정·반려·승인"),
        ("OWNER", "성과 관리", "KPI·권한·리스크"),
    ]
    for i, (role, action, desc) in enumerate(lanes):
        y = 1.85 + i * 1.08
        box(slide, 0.75, y, 2.0, 0.82, MID["panel_dark"], True)
        text(slide, role, 0.75, y, 2.0, 0.82, 12, MID["text_dark"], True, "center")
        box(slide, 2.95, y, 3.0, 0.82, "#FFFFFF", True, MID["line"])
        text(slide, action, 3.18, y + 0.08, 2.5, 0.3, 13, MID["text"], True)
        text(slide, desc, 3.18, y + 0.4, 2.5, 0.25, 9.5, MID["muted"])
        line(slide, 5.95, y + 0.4, 8.1, y + 0.4, MID["line"], 2)
        status = "자동" if i == 1 else "사람"
        fill = MID["accent"] if i == 1 else "#DDE8EA"
        box(slide, 8.12, y + 0.13, 1.05, 0.54, fill, True)
        text(slide, status, 8.12, y + 0.13, 1.05, 0.54, 11, MID["accent_text"] if i == 1 else MID["text"], True, "center")
    box(slide, 9.65, 1.85, 2.68, 4.06, "#FFF5E9", True, "#F2D3AF")
    text(slide, "승인 게이트", 9.95, 2.18, 2.05, 0.4, 15, "#8B511B", True)
    gates = ["외부 발송", "개인정보 사용", "정책·가격 판단", "대량 변경"]
    for i, item in enumerate(gates):
        text(slide, "✓", 9.95, 2.92 + i * 0.62, 0.28, 0.32, 12, "#C26A19", True)
        text(slide, item, 10.32, 2.88 + i * 0.62, 1.55, 0.38, 11, "#6F4A27", True)
    footer(slide, "AI 에이전트 업무 자동화 파일럿", MID)

    slide = blank(prs, MID["bg"])
    header(slide, "06 · ROADMAP", "8주는 진단–설계–실행–검증의 네 단계로 운영합니다", 7, MID)
    phases = [
        ("W1–2", "진단", "Process Owner", "업무·기준선"),
        ("W3–4", "설계", "Agent Squad", "프롬프트·가드레일"),
        ("W5–6", "실행", "Pilot Users", "실사용 로그"),
        ("W7–8", "검증", "PMO + Owner", "KPI·확산안"),
    ]
    for i, (week, phase, owner, deliverable) in enumerate(phases):
        x = 0.75 + i * 3.12
        color = MID["accent"] if i < 2 else MID["accent2"]
        box(slide, x, 2.0, 2.68, 3.35, "#FFFFFF", True, MID["line"])
        box(slide, x, 2.0, 2.68, 0.52, color, True)
        text(slide, week, x, 2.0, 2.68, 0.52, 11, "#FFFFFF", True, "center")
        text(slide, phase, x + 0.23, 2.82, 2.15, 0.45, 20, MID["text"], True)
        text(slide, "OWNER", x + 0.23, 3.54, 0.8, 0.25, 8.5, MID["muted"], True)
        text(slide, owner, x + 0.23, 3.82, 2.1, 0.36, 11, MID["text"], True)
        text(slide, "OUTPUT", x + 0.23, 4.38, 0.8, 0.25, 8.5, MID["muted"], True)
        text(slide, deliverable, x + 0.23, 4.66, 2.1, 0.36, 11, MID["text"], True)
    footer(slide, "AI 에이전트 업무 자동화 파일럿", MID)

    slide = blank(prs, MID["bg"])
    header(slide, "07 · DECISION", "성공을 판단할 KPI와 오늘의 승인 항목을 함께 고정합니다", 8, MID)
    kpis = [
        ("35%", "초안 시간 절감 목표", MID["accent"]),
        ("10%", "재작업률 목표", MID["accent2"]),
        ("1일", "검토 SLA 목표", "#F2A65A"),
    ]
    for i, (value, label, fill) in enumerate(kpis):
        x = 0.75 + i * 2.62
        box(slide, x, 1.92, 2.32, 1.72, "#FFFFFF", True, MID["line"])
        box(slide, x, 1.92, 0.14, 1.72, fill, False)
        text(slide, value, x + 0.32, 2.22, 1.6, 0.55, 25, MID["text"], True)
        text(slide, label, x + 0.32, 2.88, 1.7, 0.38, 10.5, MID["muted"], True)
    box(slide, 8.72, 1.92, 3.82, 3.95, MID["panel_dark"], True)
    text(slide, "승인 요청 3가지", 9.05, 2.25, 3.05, 0.45, 17, MID["text_dark"], True)
    approvals = [
        ("01", "대상 프로세스 3개"),
        ("02", "파일럿 사용자 20명"),
        ("03", "8주 운영 스쿼드"),
    ]
    for i, (num, label) in enumerate(approvals):
        y = 3.05 + i * 0.75
        circle(slide, 9.05, y, 0.42, MID["accent"])
        text(slide, num, 9.05, y, 0.42, 0.42, 9, MID["accent_text"], True, "center")
        text(slide, label, 9.62, y - 0.02, 2.2, 0.45, 12.5, MID["text_dark"], True)
    box(slide, 0.75, 4.2, 7.53, 1.67, "#DFF7F3", True)
    text(slide, "성공 게이트", 1.05, 4.5, 1.4, 0.36, 12, "#087C72", True)
    text(slide, "사용률 ≥ 60%  ·  치명 오류 0건  ·  KPI 2개 이상 개선", 1.05, 4.98, 6.65, 0.43, 15, MID["text"], True)
    footer(slide, "AI 에이전트 업무 자동화 파일럿", MID, note="모든 수치는 가정 기반 예시")

    path = PRACTICE_DIR / "PPT중간.pptx"
    prs.save(path)
    return path


def build_advanced():
    prs = new_deck()
    palette = ADV

    slide = blank(prs, palette["bg"])
    for i in range(6):
        line(slide, 8.45 + i * 0.65, 0.65, 8.45 + i * 0.65, 6.85, palette["line"], 1)
    for i in range(6):
        line(slide, 8.1, 1.0 + i * 0.95, 12.4, 1.0 + i * 0.95, palette["line"], 1)
    circle(slide, 9.2, 1.65, 2.75, palette["accent2"])
    circle(slide, 10.1, 2.55, 1.55, palette["accent"])
    pill(slide, "EXECUTIVE DECISION", 0.78, 0.72, 2.15, palette["accent"], palette["accent_text"], 9.5)
    text(slide, "AI Agent\nOperating Model 2026", 0.78, 1.58, 7.3, 1.85, 34, palette["text_dark"], True, "left", "top")
    text(slide, "90일 안에 검증하고, 12개월 안에 확장하는 운영체계", 0.82, 3.7, 6.9, 0.52, 17, palette["muted_dark"])
    text(slide, "INTERNAL SCENARIO · 2026.08", 0.82, 6.7, 4.8, 0.28, 9.5, palette["muted_dark"], True)

    slide = blank(prs, palette["bg"])
    header(slide, "01 · NORTH STAR", "12개월 목표는 ‘통제 가능한 자동화’를 조직 표준으로 만드는 것입니다", 2, palette, True)
    box(slide, 0.75, 1.92, 5.2, 3.95, palette["panel"], True, palette["line"])
    text(slide, "NORTH STAR", 1.05, 2.24, 1.5, 0.3, 10, palette["accent"], True)
    text(slide, "8,600h", 1.05, 2.82, 3.55, 0.75, 36, palette["text"], True)
    text(slide, "12개월 자동화 가능 시간", 1.08, 3.63, 3.8, 0.4, 14, palette["muted"], True)
    box(slide, 1.05, 4.47, 4.52, 0.76, "#222A37", True)
    text(slide, "목표 채택률 70%  ·  고위험 업무 자동실행 0건", 1.05, 4.47, 4.52, 0.76, 12, palette["text"], True, "center")
    text(slide, "내부 시나리오 가정", 1.08, 5.4, 2.2, 0.25, 8.5, palette["muted"])
    text(slide, "오늘 결정할 3가지", 6.55, 2.02, 4.6, 0.42, 17, palette["text"], True)
    decisions = [
        ("01", "90일 1단계 투자", "가설 검증 예산과 범위"),
        ("02", "전담 스쿼드", "업무·기술·보안 공동 오너십"),
        ("03", "데이터 접근", "최소 권한과 승인 정책"),
    ]
    for i, (num, title_value, sub) in enumerate(decisions):
        y = 2.78 + i * 1.05
        circle(slide, 6.55, y, 0.55, palette["accent"] if i == 0 else palette["accent2"])
        text(slide, num, 6.55, y, 0.55, 0.55, 10.5, palette["accent_text"], True, "center")
        text(slide, title_value, 7.35, y - 0.03, 2.6, 0.35, 14, palette["text"], True)
        text(slide, sub, 7.35, y + 0.35, 4.2, 0.3, 10, palette["muted"])
    footer(slide, "AI Agent Operating Model 2026", palette, True, note="내부 시나리오 가정")

    slide = blank(prs, palette["bg"])
    header(slide, "02 · VALUE POOL", "생산성만이 아니라 품질·속도·학습효과를 함께 측정합니다", 3, palette, True)
    pools = [
        ("8,600h", "생산성", "연간 자동화 가능 시간", palette["accent"]),
        ("−44%", "속도", "초안 리드타임 목표", palette["accent2"]),
        ("≤8%", "품질", "재작업률 목표", "#B6E36F"),
        ("70%", "학습", "월간 채택률 목표", "#FFB75E"),
    ]
    for i, (value, title_value, sub, color) in enumerate(pools):
        x = 0.72 + i * 3.12
        box(slide, x, 2.0, 2.72, 3.65, palette["panel"], True, palette["line"])
        box(slide, x, 2.0, 2.72, 0.12, color, False)
        text(slide, title_value.upper(), x + 0.28, 2.4, 2.0, 0.3, 9, color, True)
        text(slide, value, x + 0.28, 3.03, 2.15, 0.67, 27, palette["text"], True)
        text(slide, sub, x + 0.28, 3.95, 2.15, 0.7, 11.5, palette["muted"], False, "left", "top")
        box(slide, x + 0.28, 5.02, 2.1, 0.32, "#252C38", True)
        text(slide, "BASELINE → TARGET", x + 0.28, 5.02, 2.1, 0.32, 8.3, palette["muted"], True, "center")
    footer(slide, "AI Agent Operating Model 2026", palette, True, note="모든 수치는 내부 시나리오 가정")

    slide = blank(prs, palette["bg"])
    header(slide, "03 · MATURITY", "Assist에서 Transform으로 갈수록 자동화보다 운영 통제가 중요해집니다", 4, palette, True)
    stages = [
        ("01", "ASSIST", "개인 초안", "현재"),
        ("02", "AUTOMATE", "규칙 기반 실행", "90일"),
        ("03", "ORCHESTRATE", "업무 간 연결", "6개월"),
        ("04", "TRANSFORM", "운영모델 재설계", "12개월"),
    ]
    for i, (num, name, desc, timing) in enumerate(stages):
        x = 0.72 + i * 3.12
        y = 4.72 - i * 0.62
        h = 1.22 + i * 0.62
        fill = "#1D2632" if i < 2 else "#242437"
        box(slide, x, y, 2.72, h, fill, True, palette["line"])
        text(slide, num, x + 0.25, y + 0.22, 0.42, 0.32, 9, palette["accent"] if i < 2 else palette["accent2"], True)
        text(slide, name, x + 0.25, y + 0.68, 2.1, 0.38, 15, palette["text"], True)
        text(slide, desc, x + 0.25, y + 1.15, 2.1, 0.35, 10, palette["muted"])
        pill(slide, timing, x + 0.25, y + h - 0.58, 0.92, "#303847", palette["text"], 8.5)
    text(slide, "자율성", 0.78, 6.33, 0.85, 0.3, 9, palette["muted"], True)
    line(slide, 1.65, 6.49, 12.3, 6.49, palette["accent"], 3)
    text(slide, "높음 →", 11.5, 6.33, 0.8, 0.3, 9, palette["accent"], True, "right")
    footer(slide, "AI Agent Operating Model 2026", palette, True)

    slide = blank(prs, palette["bg"])
    header(slide, "04 · TARGET ARCHITECTURE", "5개 계층을 분리해야 모델이 바뀌어도 운영체계는 유지됩니다", 5, palette, True)
    layers = [
        ("EXPERIENCE", "Chat · Portal · Workflow", "#25313D"),
        ("ORCHESTRATION", "Router · Planner · Memory", "#233044"),
        ("SKILLS", "Proposal · VOC · Report · QA", "#292747"),
        ("DATA", "Knowledge · CRM · Documents", "#222A37"),
        ("GOVERNANCE", "Policy · Approval · Audit · Kill Switch", "#2F2A25"),
    ]
    for i, (name, items, fill) in enumerate(layers):
        y = 1.82 + i * 0.86
        box(slide, 0.82, y, 11.7, 0.66, fill, True, palette["line"])
        text(slide, name, 1.08, y, 2.3, 0.66, 11.5, palette["accent"] if i < 2 else palette["accent2"], True)
        text(slide, items, 3.65, y, 7.95, 0.66, 12, palette["text"], True)
        text(slide, f"L{i+1}", 11.75, y, 0.45, 0.66, 9.5, palette["muted"], True, "right")
    box(slide, 2.3, 6.32, 8.75, 0.48, "#202833", True)
    text(slide, "원칙  ·  최소 권한  /  추적 가능  /  사람 승인  /  즉시 중단 가능", 2.3, 6.32, 8.75, 0.48, 10.5, palette["text"], True, "center")
    footer(slide, "AI Agent Operating Model 2026", palette, True)

    slide = blank(prs, palette["bg"])
    header(slide, "05 · TRUST BY DESIGN", "위험 등급에 따라 자동화 범위와 승인 게이트를 다르게 설계합니다", 6, palette, True)
    risks = [
        ("L1", "LOW", "요약·분류", "자동 실행", palette["accent"]),
        ("L2", "MEDIUM", "내부 초안", "사후 표본 검토", palette["accent2"]),
        ("L3", "HIGH", "외부 발송", "사전 사람 승인", "#FFB75E"),
        ("L4", "CRITICAL", "정책·금전", "자동 실행 금지", "#FF6B73"),
    ]
    for i, (level, name, scope, control, fill) in enumerate(risks):
        y = 1.86 + i * 1.08
        box(slide, 0.78, y, 11.78, 0.84, palette["panel"], True, palette["line"])
        box(slide, 0.78, y, 0.9, 0.84, fill, True)
        text(slide, level, 0.78, y, 0.9, 0.84, 12, palette["accent_text"], True, "center")
        text(slide, name, 2.0, y, 1.55, 0.84, 12, palette["text"], True)
        text(slide, scope, 4.0, y, 2.2, 0.84, 12, palette["muted"], True)
        text(slide, control, 7.0, y, 2.65, 0.84, 12.5, fill, True)
        text(slide, "AUDIT LOG", 10.45, y, 1.25, 0.84, 8.5, palette["muted"], True, "center")
    box(slide, 2.15, 6.33, 9.0, 0.48, "#25242F", True)
    text(slide, "Kill Switch는 모든 등급에 공통 적용하고, 권한·프롬프트·출력 이력을 보존합니다.", 2.15, 6.33, 9.0, 0.48, 10.5, palette["text"], True, "center")
    footer(slide, "AI Agent Operating Model 2026", palette, True)

    slide = blank(prs, palette["bg"])
    header(slide, "06 · PORTFOLIO", "가치와 실행 가능성이 모두 높은 두 업무를 1단계에 배치합니다", 7, palette, True)
    box(slide, 1.1, 1.78, 10.5, 4.78, "#171C25", False, palette["line"])
    line(slide, 6.35, 1.78, 6.35, 6.56, palette["line"], 2)
    line(slide, 1.1, 4.17, 11.6, 4.17, palette["line"], 2)
    text(slide, "VALUE ↑", 0.45, 1.98, 0.55, 0.5, 8.5, palette["muted"], True, "center")
    text(slide, "FEASIBILITY →", 9.65, 6.64, 1.95, 0.3, 8.5, palette["muted"], True, "right")
    bubbles = [
        (3.0, 2.42, 1.15, "주간 보고", palette["accent"]),
        (4.55, 3.17, 1.0, "VOC 분류", palette["accent"]),
        (7.25, 2.35, 1.28, "제안서 초안", palette["accent2"]),
        (8.9, 4.78, 1.12, "캠페인 QA", "#FFB75E"),
        (5.0, 5.05, 0.98, "지식 검색", "#667085"),
    ]
    for x, y, d, label, fill in bubbles:
        circle(slide, x, y, d, fill)
        text(slide, label, x - 0.12, y + d / 2 - 0.17, d + 0.24, 0.34, 9.5, "#071316" if fill == palette["accent"] else "#FFFFFF", True, "center")
    pill(slide, "WAVE 1", 1.36, 2.05, 1.12, "#243E45", palette["accent"], 9)
    footer(slide, "AI Agent Operating Model 2026", palette, True, note="우선순위는 내부 시나리오 가정")

    slide = blank(prs, palette["bg"])
    header(slide, "07 · BUSINESS CASE", "잠재 시간가치는 3.612억원, 1년 가치/투자 비율은 1.72×입니다", 8, palette, True)
    cards = [
        ("2.1억원", "1년 총투자", palette["accent2"]),
        ("3.612억원", "잠재 연간 시간가치", palette["accent"]),
        ("1.72×", "단순 가치/투자", "#B6E36F"),
    ]
    for i, (value, label, fill) in enumerate(cards):
        x = 0.75 + i * 3.05
        box(slide, x, 1.88, 2.7, 1.62, palette["panel"], True, palette["line"])
        box(slide, x, 1.88, 0.11, 1.62, fill, False)
        text(slide, value, x + 0.28, 2.18, 2.05, 0.55, 24, palette["text"], True)
        text(slide, label, x + 0.28, 2.86, 2.05, 0.3, 10.5, palette["muted"], True)
    box(slide, 10.1, 1.88, 2.45, 1.62, "#25242F", True, palette["line"])
    text(slide, "산식", 10.38, 2.17, 0.8, 0.3, 9.5, palette["accent2"], True)
    text(slide, "8,600h × 42,000원", 10.38, 2.64, 1.78, 0.4, 12.5, palette["text"], True)
    scenarios = [("50%", 1.81, "#667085"), ("70%", 2.53, palette["accent2"]), ("100%", 3.61, palette["accent"])]
    text(slide, "실현률 민감도 · 시간가치(억원)", 0.78, 4.06, 4.0, 0.36, 12, palette["text"], True)
    for i, (rate, value, fill) in enumerate(scenarios):
        y = 4.68 + i * 0.58
        text(slide, rate, 0.82, y, 0.62, 0.3, 10, palette["muted"], True)
        box(slide, 1.55, y + 0.04, value * 2.35, 0.24, fill, True)
        text(slide, f"{value:.2f}", 10.28, y, 0.75, 0.3, 10, palette["text"], True, "right")
    box(slide, 10.15, 4.35, 2.4, 1.72, "#222A37", True)
    text(slide, "주의", 10.45, 4.67, 0.72, 0.3, 9.5, "#FFB75E", True)
    text(slide, "실제 현금 절감이 아닌\n잠재 시간가치입니다.", 10.45, 5.08, 1.75, 0.68, 11, palette["text"], True, "left", "top")
    footer(slide, "AI Agent Operating Model 2026", palette, True, note="내부 시나리오 가정 · 세전 단순 비교")

    slide = blank(prs, palette["bg"])
    header(slide, "08 · 90-DAY PLAN", "각 단계의 산출물과 Scale Gate를 먼저 고정합니다", 9, palette, True)
    phases = [
        ("D01–20", "DISCOVER", "업무·리스크 기준선", palette["accent"]),
        ("D21–45", "BUILD", "스킬·가드레일·로그", palette["accent2"]),
        ("D46–75", "PILOT", "실사용·KPI 측정", "#B6E36F"),
        ("D76–90", "SCALE GATE", "확산/중단 의사결정", "#FFB75E"),
    ]
    line(slide, 1.0, 3.42, 12.1, 3.42, palette["line"], 4)
    for i, (days, name, output, fill) in enumerate(phases):
        x = 0.72 + i * 3.12
        circle(slide, x + 0.98, 3.02, 0.82, fill)
        text(slide, f"{i+1}", x + 0.98, 3.02, 0.82, 0.82, 15, palette["accent_text"], True, "center")
        text(slide, days, x, 1.95, 2.8, 0.32, 10, fill, True, "center")
        text(slide, name, x, 2.35, 2.8, 0.42, 14.5, palette["text"], True, "center")
        box(slide, x, 4.2, 2.8, 1.25, palette["panel"], True, palette["line"])
        text(slide, output, x + 0.2, 4.47, 2.4, 0.48, 11.5, palette["text"], True, "center")
        text(slide, ["CPO", "Tech Lead", "Pilot Owner", "SteerCo"][i], x, 5.65, 2.8, 0.28, 9, palette["muted"], True, "center")
    footer(slide, "AI Agent Operating Model 2026", palette, True)

    slide = blank(prs, palette["bg"])
    header(slide, "09 · DECISION", "90일 투자는 세 가지 승인과 네 개 성공 게이트로 통제합니다", 10, palette, True)
    text(slide, "승인 요청", 0.78, 1.9, 3.0, 0.42, 16, palette["text"], True)
    approvals = [
        ("01", "1단계 투자 2.1억원 한도"),
        ("02", "전담 스쿼드 6명 구성"),
        ("03", "최소 권한 데이터 접근"),
    ]
    for i, (num, label) in enumerate(approvals):
        y = 2.55 + i * 0.9
        box(slide, 0.78, y, 5.25, 0.68, palette["panel"], True, palette["line"])
        text(slide, num, 1.02, y, 0.48, 0.68, 10.5, palette["accent"], True, "center")
        text(slide, label, 1.65, y, 3.9, 0.68, 12.5, palette["text"], True)
    text(slide, "Scale Gate", 6.65, 1.9, 3.0, 0.42, 16, palette["text"], True)
    gates = [
        ("G1", "채택률 ≥ 60%"),
        ("G2", "KPI 2개 이상 개선"),
        ("G3", "고위험 자동 실행 0건"),
        ("G4", "감사 로그 완전성 100%"),
    ]
    for i, (num, label) in enumerate(gates):
        x = 6.65 + (i % 2) * 2.88
        y = 2.55 + (i // 2) * 1.32
        box(slide, x, y, 2.58, 1.0, "#202A33", True, palette["line"])
        text(slide, num, x + 0.2, y + 0.18, 0.55, 0.3, 9.5, "#B6E36F", True)
        text(slide, label, x + 0.2, y + 0.5, 2.1, 0.3, 11, palette["text"], True)
    box(slide, 0.78, 5.75, 11.45, 0.68, palette["accent"], True)
    text(slide, "NEXT 72 HOURS  ·  오너 확정 → 킥오프 → 데이터 접근 워크숍", 0.78, 5.75, 11.45, 0.68, 14, palette["accent_text"], True, "center")
    footer(slide, "AI Agent Operating Model 2026", palette, True, note="경영진 의사결정용 · 내부 시나리오")

    path = PRACTICE_DIR / "ppt심화.pptx"
    prs.save(path)
    return path


def validate_deck(path: Path, expected_slides: int):
    with zipfile.ZipFile(path) as archive:
        bad_member = archive.testzip()
        if bad_member:
            raise RuntimeError(f"{path.name}: corrupt OOXML member {bad_member}")
        required = {"[Content_Types].xml", "ppt/presentation.xml"}
        missing_required = required.difference(archive.namelist())
        if missing_required:
            raise RuntimeError(f"{path.name}: missing {sorted(missing_required)}")

    prs = Presentation(path)
    if len(prs.slides) != expected_slides:
        raise RuntimeError(
            f"{path.name}: expected {expected_slides} slides, got {len(prs.slides)}"
        )
    ratio = prs.slide_width / prs.slide_height
    if abs(ratio - (16 / 9)) > 0.01:
        raise RuntimeError(f"{path.name}: invalid aspect ratio {ratio}")

    out_of_bounds = []
    unresolved = []
    all_text = []
    for slide_no, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if (
                shape.left < 0
                or shape.top < 0
                or shape.left + shape.width > prs.slide_width + 5000
                or shape.top + shape.height > prs.slide_height + 5000
            ):
                out_of_bounds.append((slide_no, shape.name))
            if hasattr(shape, "text"):
                value = shape.text or ""
                all_text.append(value)
                if re.search(r"TODO|\{\{.+?\}\}|PLACEHOLDER", value, re.I):
                    unresolved.append((slide_no, value))
    if out_of_bounds:
        raise RuntimeError(f"{path.name}: out-of-bounds shapes {out_of_bounds}")
    if unresolved:
        raise RuntimeError(f"{path.name}: unresolved text {unresolved}")
    if not "".join(all_text).strip():
        raise RuntimeError(f"{path.name}: no text found")
    return {
        "file": path.name,
        "slides": len(prs.slides),
        "ratio": round(ratio, 4),
        "shapes": sum(len(slide.shapes) for slide in prs.slides),
    }


def main():
    outputs = [build_basic(), build_mid(), build_advanced()]
    expected = {
        "PPT기초.pptx": 6,
        "PPT중간.pptx": 8,
        "ppt심화.pptx": 10,
    }
    print("PPT PRACTICE BUILD COMPLETE")
    for path in outputs:
        result = validate_deck(path, expected[path.name])
        print(
            f"[PASS] {result['file']} | slides={result['slides']} "
            f"| ratio={result['ratio']} | shapes={result['shapes']}"
        )


if __name__ == "__main__":
    main()
