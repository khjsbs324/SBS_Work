import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parents[1] / "pptdeps"))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Pt


src, out = map(Path, sys.argv[1:3])
prs = Presentation(src)

FONT = "Noto Sans KR"
NAVY = RGBColor(6, 70, 130)
BLUE = RGBColor(10, 79, 153)
DARK = RGBColor(10, 53, 106)
PALE = RGBColor(238, 247, 255)
PALE2 = RGBColor(221, 240, 255)
BORDER = RGBColor(185, 217, 245)
TEXT = RGBColor(30, 43, 56)
WHITE = RGBColor(255, 255, 255)


def title_of(slide):
    values = []
    for shape in slide.shapes:
        if (
            getattr(shape, "has_text_frame", False)
            and shape.text.strip()
            and shape.top.pt < 90
            and shape.left.pt < 850
            and not shape.text.strip().startswith("CHAPTER")
        ):
            values.append(shape.text.strip())
    return max(values, key=len) if values else ""


def find(title):
    return next(slide for slide in prs.slides if title_of(slide) == title)


def remove(shape):
    shape._element.getparent().remove(shape._element)


def clear_body(slide):
    for shape in list(slide.shapes):
        if shape.top.pt > 82:
            remove(shape)


def text(slide, x, y, w, h, value, size=12, bold=False, color=TEXT, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(Pt(x), Pt(y), Pt(w), Pt(h))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = Pt(1)
    frame.margin_top = frame.margin_bottom = Pt(1)
    frame.vertical_anchor = MSO_ANCHOR.TOP
    for index, line in enumerate(str(value).split("\n")):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.alignment = align
        paragraph.space_after = Pt(1)
        run = paragraph.add_run()
        run.text = line
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return shape


def box(slide, x, y, w, h, value="", fill=WHITE, line=BORDER, size=12,
        bold=False, color=TEXT, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Pt(x), Pt(y), Pt(w), Pt(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1)
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = Pt(14)
    frame.margin_top = frame.margin_bottom = Pt(7)
    frame.vertical_anchor = MSO_ANCHOR.TOP
    for index, line_value in enumerate(str(value).split("\n")):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.alignment = align
        paragraph.space_after = Pt(1)
        run = paragraph.add_run()
        run.text = line_value
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return shape


def exact_page(slide, content, size=12):
    clear_body(slide)
    box(slide, 42, 98, 875, 398, content, WHITE, BORDER, size, False, TEXT)


def native_table(slide, x, y, w, h, headers, rows, size=11.5, widths=None):
    shape = slide.shapes.add_table(len(rows) + 1, len(headers), Pt(x), Pt(y), Pt(w), Pt(h))
    table = shape.table
    if widths:
        for index, width in enumerate(widths):
            table.columns[index].width = Pt(width)
    for row_index, row in enumerate([headers] + rows):
        for column_index, value in enumerate(row):
            cell = table.cell(row_index, column_index)
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY if row_index == 0 else (PALE if row_index % 2 else WHITE)
            cell.margin_left = cell.margin_right = Pt(4)
            cell.margin_top = cell.margin_bottom = Pt(3)
            frame = cell.text_frame
            frame.clear()
            frame.word_wrap = True
            frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            paragraph = frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.CENTER
            run = paragraph.add_run()
            run.text = str(value)
            run.font.name = FONT
            run.font.size = Pt(size)
            run.font.bold = row_index == 0 or column_index == 0
            run.font.color.rgb = WHITE if row_index == 0 else TEXT
    return shape


def set_title(slide, value, badge=None):
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and shape.top.pt < 90 and shape.left.pt < 850:
            if shape.text.strip() and not shape.text.strip().startswith("CHAPTER"):
                shape.text_frame.paragraphs[0].runs[0].text = value
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = FONT
                        run.font.size = Pt(22 if len(value) > 18 else 25)
                        run.font.bold = True
                        run.font.color.rgb = DARK
                break
    if badge is not None:
        for shape in slide.shapes:
            if shape.name == "PageBadge" or (getattr(shape, "has_text_frame", False) and shape.top.pt < 60 and shape.left.pt > 850):
                shape.text = badge
                for paragraph in shape.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.CENTER
                    for run in paragraph.runs:
                        run.font.name = FONT
                        run.font.size = Pt(10.5)
                        run.font.bold = True
                        run.font.color.rgb = BLUE
                shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                break


def move_after(new_slide, anchor):
    ids = prs.slides._sldIdLst
    new_node = next(node for node in ids if int(node.id) == new_slide.slide_id)
    anchor_node = next(node for node in ids if int(node.id) == anchor.slide_id)
    ids.remove(new_node)
    ids.insert(list(ids).index(anchor_node) + 1, new_node)


def header(slide, chapter, title_value, badge):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    text(slide, 42, 9, 520, 14, chapter, 8.5, True, RGBColor(138, 166, 199))
    title_shape = text(slide, 42, 29, 800, 42, title_value, 22 if len(title_value) > 18 else 25, True, DARK)
    title_shape.name = "Title"
    badge_shape = box(slide, 878, 28, 39, 28, badge, PALE2, PALE2, 10.5, True, BLUE, PP_ALIGN.CENTER)
    badge_shape.name = "PageBadge"
    badge_shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Pt(42), Pt(79), Pt(875), Pt(1.4))
    rule.fill.solid()
    rule.fill.fore_color.rgb = BLUE
    rule.line.fill.background()


def new_after(anchor, chapter, title_value, badge):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    move_after(slide, anchor)
    header(slide, chapter, title_value, badge)
    return slide


# 06: source flowchart wording, including the source's original spacing and phrasing.
exact_page(find("국비 문의의 흐름 · 원문 1"), """“국비되나요?”
카드발급 받으셨어요?

사용해보셨을까요? → 네
사용해봤던 사람들은 국비의 단점 (피드백안되는점, 일을병행할수없는점, 수업기간이 긴점) 을 잘 알고있기때문에, 이 부분만 살짝만 건들어준다면, 국비타파가 쉽다

ex) 국비수업을 들어봐서 아시겠지만, 수업시간이 고정이 되어있어서 지금상황에서는 활용이 안되실꺼에요. 그러면 본원자체에 국비카드 소지자 분들은 저희가 따로 20%~40% 정도 할인적용해서 자부담금내시고 수업 들으실수있으세요""", 12.5)

exact_page(find("국비 문의의 흐름 · 원문 2"), """사용해보셨을까요? → 아니요
앞서 말했듯이 카드는 발급을 받았지만 활용하지 못하는 경우가 많다. 이럴경우 문의자의 수업 목적, 시간대, 기간을 파악하여 반대로 국비카드를 지금도 역시나 사용하지 못하신다라는 인식을 심어줄것!

문의자 상황 : 알바를 하면서 국비카드 사용되나요? / 직장을 다니고있는데 주말반이나, 저녁반으로 가능한가요?
ex) 국비수업은 9시 ~ 16시까지 수업에만 활용가능하신데 그 시간 수업 들어오실수있으세요?
=> 직장인들 참여불가

국비수업은 취업위주 수업이기때문에 저녁반이나, 주말반은 수업이 없습니다 => 직장인들 참여불가""", 12)

exact_page(find("국비 문의의 흐름 · 원문 3"), """발급신청은 하셨어요? → 네
내배카로 신청하셨을까요, 국민취업지원제도로 신청하셨을까요?
먼저는 국비종류 파악이 먼저! 내배카로 신청했을경우 90%이상 그럴경우, 국취제가 좋고, 활용하실수있는 범위가 넓다고 하여서 국취제로 신청하게끔 진행
☆ 국취제의 장점 어필

ex) 내배카 활용하실려면 지금 일 그만두시고 오전부터 낮까지 6개월동안 수업 들어가야하는데 다른일을 병행을 할수없기때문에, 국취제를 활용해서 이직or취업준비하시면서 간단한 알바는 다른 취업준비를 하실수있기때문에 국취제로 신청하면됩니다~

발급신청은 하셨어요? → 아니요
국비의 종류가 많기때문에, 이야기를 통해서 본인에게 맞는 국비나라지원을 파악해줄께요

★ 제일쉬운 컨택 및 상담, 등록의 상황""", 11.7)


# 07: retain the unusual source word order instead of correcting it.
exact_page(find("국비 문의자 유형 · 원문"), """국비문의자 유형

내일배움카드 1회이상 사용자 15%
내일배움카드 소지자 40%
내일배움카드 미소지자 45%

“국비” 겁먹지말고, 자신있게 컨택하자

“국비지원제도, 국비가능한가요?”
우리가 어떠한 물건을 살때 “혹시 현금으로 하면 할인되나요?”라는 질문과 동일!
만약, 현금으로 할인이 될 경우에는 안살것인가?
NO! 구매목적이 있었다면, 현금과 카드는 의미는 없다
그러나, 소비자 입장에서는 좀더 할인되는 방법을 찾고 싶을뿐
없다면 그대로 구매한다

실제로 문의준 DB를 컨택을 해보면, 내배카 or 국취제를 아는 정확히 문의자는 없다.
배움에 목적이 있다면, 국비는 크게 문제될건 없다.""", 12.5)


# 10: exact DB values and exact consultation wording.
slide = find("페이스북 DB 컨택 스피치 · 원문 DB")
clear_body(slide)
box(slide, 42, 100, 875, 178, """지점 : 대전
이름 : 함승범(4325163)
상담결과 : 완납 6,030,000
핸드폰 : 010-4802-1595
문의과목 : 드로잉
DB등록일자 기준 : 2025-08-01 오후 5:00:13
담당자 : [대전추천2-1]김호영2
학과 : 아트웍
출처 : 디비추천(함승보 010-8543-1595)""", PALE, BORDER, 11.8)
box(slide, 42, 296, 875, 178, """지점 : 대전
이름 : 함승보(4324482)
상담결과 : 완납 6,930,000
핸드폰 : 010-8543-1595
문의과목 : 영상편집
DB등록일자 기준 : 2025-08-01 오전 9:17:45
담당자 : [대전페북2-1]김호영2
학과 : 모션/영상
출처 : [페이스북Meta]_(대전) 국비 영상편집""", WHITE, BORDER, 11.8)

exact_page(find("페이스북 DB 컨택 스피치 · 원문 상담"), """상담예정 08.01 ★16:00
*컨택자 : 김호영
*본인확인 및 문의여부 체크 : 본인
*배워보시려는 계기 및 수강 희망 시기 : 취미 및 취업
*분야 경험 및 수강 경험 여부 : 프리미어 다뤄봄
*현재 거주지 및 평소 스케줄 어떤지 (가용시간체크) : 태평동
*상담 명분 제시 및 상담 일정 확정 : 국비체크
*이외 상담내용 :
미디어콘텐츠학과 1학년 다니다가 자퇴했다고 함 -> 적성은 맞으나 학교가 맞지 않음
프리미어는 조금 다룰 수 있고 나머진 해보기만 했었다고 하였음
군대는 가지 않았지만 아직 예정은 없다고 하였음 -> 취업, 취미 둘 다 생각중이라고 함
2D 영상편집 외에도 다른 과정들이 많다고 안내해드렸고 국비 과정 설명해드리기로 하였음
2시에 둔산동 헬스 레슨이 있어서 온다고 했으며 오전 수업은 선호하지 않음
9월부터 수업 희망하고 오전 수업 희망하지 않아 영상계좌제 비교하면서 설명할 예정
**수강료 / 시간표에 대한 나의 응대 : X
*같이 방문하실 지인 / 기타 특이사항 : X
*통화시간 : 6분""", 11.3)


# 11: exact record and KakaoTalk transcript.
exact_page(find("실업급여 수령자 컨택 스피치 · 원문 상담"), """*컨택자 : 김호영
*본인확인 및 문의여부체크 : 본인
*배워보시려는 계기 및 수강 희망 시기 : 취업
*분야 경험 및 수강 경험 여부 : X
*현재 거주지 및 평소 스케줄 어떻게 (가용시간체크) : 대전
*상담 명분 제시 및 상담 일정 확정 : 익일 2시
*이외 상담내용 :
유선으로 연결되었고 영상이 아닌 웹디자인 국비 궁금하다고 함
1년 전에 뷰티과정 1년 국비로 수강했고 추가로 금액도 지불했다고 함
전에 근무하면서 가볍게 디자인을 경험할 일이 있었는데 잘 맞는 것 같아서 해보려고 함
실업급여 수령 중인데 이제 2회차를 받는다고 함 -> 매월 5일 기준이고 6개월 받는다고 함
실업급여랑 국비랑 병행 안되는지 여쭤봐서 실업급여 받기 위해서는 구직활동 필요하니까 수업 듣는거 좋다고 안내
내배카 잔액이랑 체크해드리기 위해 익일 상담 도와드리기로 함
**수강료 / 시간표에 대한 나의 응대 : X
*같이 방문하실 지인 / 기타 특이사항 : X
*통화시간 : 7분""", 11.5)

exact_page(find("실업급여 수령자 컨택 스피치 · 원문 대화"), """★박지유(오피스)
멘토님 저 실업급여 제출서류가 필요한데 수강 신청서랑 mos, 지금 듣고있는 컴활2급 출석부 좀 보내주세요~!
오후 1:48

몇 일부터 몇 일까지 필요하실까요?
오후 1:48

★박지유(오피스)
Mos랑 컴활 처음 수강날부터 내일까지 막날까지요!
오후 1:49""", 14)


# 20: keep the source notice as written, without editorial punctuation changes.
exact_page(find("졸업예정자 상담 · 원문"), """국취 신청기준

@국민취업지원제도 참여제한: 주30시간이상 고용보험가입자/실업급여수급자(2유형은 수급종료후다음날부터 참여가능)/기초생계수급자(본인, 가족포함)(단,교육, 의료, 주거급여 수급자는 가능), 26년 졸업예정자중 2월졸업자(25.07.01부터), 26년 8월졸업예정자(25.11.01부터)/진학예정자/인문계고 졸업예정자는 졸업년도1.1부터(특성화고는 3학년 당해연도7.1부터/1유형 신청자 본인 소득이 1,435,208(1인 중위소득 60%이상) 이상시 1유형신청불가/사업자등록증 보유자는 월 250만원 미만 소득 또는 월 1,250만원 미만 매출 발생

※ 1유형 소득감액기준: 2025년 기준금액 1,435,208원 (구촉 월지급액 50~70만원인경우 935,208이상 소득발생시 부지급/구촉 월지급액 80만원인경우 80만원이상 소득발생시 부지급)""", 11.8)


# 21: reproduce all source table columns, labels, figures, and blank cells.
slide = find("견적서 활용 · 3D그래픽학과 1")
clear_body(slide)
native_table(slide, 42, 105, 875, 330, ["", "과목명", "COURSE", "등록금"], [
    ["9월", "프리미어 프로", "1개월", "400,000"],
    ["10월", "에프터이펙트 기초", "1개월", "450,000"],
    ["11월", "에프터이펙트 활용", "1개월", "450,000"],
    ["12월", "에프터이펙트 심화", "1개월", "450,000"],
    ["9월", "블렌더 기초", "1개월", "500,000"],
    ["10월", "블렌더 활용", "1개월", "500,000"],
    ["11월", "3D 에셋&지브러쉬 기초", "1개월", "600,000"],
], 12, [90, 500, 120, 165])

slide = find("견적서 활용 · 3D그래픽학과 2")
clear_body(slide)
native_table(slide, 42, 105, 875, 330, ["", "과목명", "COURSE", "등록금"], [
    ["12월", "3D 에셋&지브러쉬 활용", "1개월", "600,000"],
    ["9월", "렌더&라이팅 기초", "1개월", "600,000"],
    ["10월", "렌더&라이팅 활용", "1개월", "600,000"],
    ["11월", "리깅&애니메이션 기초", "1개월", "600,000"],
    ["12월", "리깅&애니메이션 활용", "1개월", "600,000"],
    ["1월", "룩뎁&언리얼", "1개월", "600,000"],
    ["12월", "포트폴리오 (무료과목)", "6개월", "3,600,000"],
], 12, [90, 500, 120, 165])
box(slide, 142, 458, 675, 30, "총 등록금 10,550,000", PALE2, BORDER, 13, True, DARK, PP_ALIGN.CENTER)

slide = find("견적서 활용 · 기계학과 1")
clear_body(slide)
native_table(slide, 42, 105, 875, 330, ["", "과목명", "COURSE", "등록금"], [
    ["4월", "캐드 기초", "1개월", "350,000"],
    ["5월", "캐드 활용", "1개월", "350,000"],
    ["6월", "인벤터", "1개월", "400,000"],
    ["7월", "퓨전 360 기초", "1개월", "400,000"],
    ["8월", "퓨전 360 활용", "1개월", "400,000"],
    ["8월", "전산응용기계제도기능사", "1개월", "450,000"],
    ["9월", "사무 기초", "1개월", "300,000"],
], 12, [90, 500, 120, 165])

slide = find("견적서 활용 · 기계학과 2")
clear_body(slide)
native_table(slide, 42, 105, 875, 330, ["", "과목명", "COURSE", "등록금"], [
    ["10월", "컴퓨터활용능력 1급", "1개월", "350,000"],
    ["5월", "AI 프롬프트 기초 (무료교육)", "1개월", "600,000"],
    ["6월", "AI 프롬프트 활용 (무료교육)", "1개월", "600,000"],
    ["5월", "AI 에이전트 기초 (무료교육)", "1개월", "600,000"],
    ["6월", "AI 에이전트 활용 (무료교육)", "1개월", "600,000"],
    ["", "재수강", "1개월", "-"],
    ["", "재수강", "1개월", "-"],
], 12, [90, 500, 120, 165])
box(slide, 142, 458, 675, 30, "총 등록금 5,400,000", PALE2, BORDER, 13, True, DARK, PP_ALIGN.CENTER)

slide = find("견적서 활용 · 3D그래픽 장학지원")
clear_body(slide)
native_table(slide, 42, 112, 875, 260,
             ["세부장학지원 내역", "장학지원율(%)", "지원금액", "등록금"], [
                 ["AIX페이 (~08.12)", "", "3,300,000", "7,250,000"],
                 ["온라인결합", "", "1,000,000", "6,250,000"],
                 ["당일등록", "", "355,000", "5,895,000"],
                 ["국민취업지원제도 1유형", "", "3,600,000", "2,295,000"],
             ], 11.5, [390, 150, 160, 175])
box(slide, 92, 410, 775, 62,
    "최초 납부해야할 등록금\n3,600,000(국취제) + 2,295,000(자부담) + 990,000(온라인)",
    PALE2, BORDER, 14, True, DARK, PP_ALIGN.CENTER)

slide = find("견적서 활용 · 기계학과 장학지원")
clear_body(slide)
native_table(slide, 42, 112, 875, 230,
             ["세부장학지원 내역", "장학지원율(%)", "지원금액", "등록금"], [
                 ["새학기 프로모션 (~04.13)", "", "1,800,000", "3,600,000"],
                 ["온라인 결합", "", "600,000", "3,000,000"],
                 ["국민취업지원제도 1유형", "", "3,600,000", "-"],
             ], 11.5, [390, 150, 160, 175])
box(slide, 92, 390, 775, 62,
    "총 납부해야할 등록금\n3,000,000(국취제) + 500,000(온라인)",
    PALE2, BORDER, 14, True, DARK, PP_ALIGN.CENTER)


# 23: exact KakaoTalk names, dates, messages, and source typos.
exact_page(find("카카오톡 활용 · 원문 1"), """https://blog.naver.com/pj3503/222485497313
주민등록 주소지 이전변경
주민등록 주소지 이전 변경 방법 인터넷으로 할수있습니다.
blog.naver.com
오후 3:28

국민취업지원...신청 방법.pdf
유효기간 ~2022. 8. 9. / 용량 387 KB
오후 3:29

중원님 주소지 옮기고 국민취업지원제도 1유형으로 신청하시면 되세요~ 신청하시다가 헷갈리는 부분 있으시면 연락주세요
오후 3:43

★최중원(코딩)
주소지 이전 신청은 했는데
오후 4:22

★최중원(코딩)
세대주가 확인을 해야하는데 이걸 인증을 하려면 누나가 시간이 아직 안되어서 오늘안에 국비지원신청은 못할것 같은데 선결제는 내일 가능한데 천천히 신청되어도 괜찮은가여?
오후 4:23""", 11.5)

exact_page(find("카카오톡 활용 · 원문 2"), """★권준혁(실내건축)
통화 끝났고 등록 할게요
오후 2:58

https://blog.naver.com/pj3503/222485497313
오후 3:34

주민등록 주소지 이전변경
주민등록 주소지 이전 변경 방법 인터넷으로 할수있습니다.
blog.naver.com

국민취업지원...신청 방법.pdf
유효기간 ~2023. 5. 11. / 용량 387 KB

주소지 먼저 옮기고 국민취업지원제도 1유형 선발형(청년)으로 신청해주세요~
오후 3:34""", 12.5)

exact_page(find("카카오톡 활용 · 원문 3"), """https://blog.naver.com/pj3503/222485497313
주민등록 주소지 이전변경
주민등록 주소지 이전 변경 방법 인터넷으로 할수있습니다.
blog.naver.com

국민취업지원...신청 방법.pdf
유효기간 ~2023. 5. 29. / 용량 387 KB

내용 확인하고 신청하다가 헷갈리시면 말씀해주세요!
오후 6:57

★김도형(기계)
네 해볼게요
오후 7:07

★김도형(기계)
신청은 해놨어요
오후 7:29""", 13)


# 24: source chat and deposit screens; blank source cells stay blank.
slide = find("기존 수강생 수급내역 · 원문")
clear_body(slide)
box(slide, 42, 98, 875, 180, """★김효령(마야)
멘토님
이거 들어온 것 맞나요??
오후 5:53

★김효령(마야)
3월 수당이욤..!
오후 5:54""", WHITE, BORDER, 14)
native_table(slide, 92, 292, 775, 170, ["구분", "일자·시간", "입금액", "잔액"], [
    ["광주고용청", "", "500,000원", "789,353원"],
    ["국민취업지원", "3월 22일 16:15", "500,000원", "585,267원"],
    ["훈련비고용부", "5월 17일 16:59", "200,000원", "380,176원"],
], 12, [190, 200, 180, 205])


# 25: both evidence records, with every visible field retained verbatim.
exact_page(find("국취제 과정 확대 · 원문 1"), """이름 : 이수진
휴대폰 : 010-5095-1761
나이 : 20대
성별 : ○남 ◉여
전공유무 : 비전공자
공부목적 : □ 진학 ■ 취업 □ 창업 □ 취미
문의과목 : 구직자, 전산세무 회계 / 오피스 자격증
희망교육장 : 대전
당직여부 : ■ 당직
출처경로 : [당] www.sbsart.com 국비지원과정조회
출처담당 : 김홍구5 (1588-5530)
상담신청일 : 2023-02-27 오후 6:37:02
희망수강일정 : 월 / 희망시간 / 평일/주말
담당매니저 : [대전당직2-1] 김호영2
학과구분 : 자격증
간단메모 : *DB 담당자 혹은 광고자를 입력해주세요.

컨택결과 : 재컨택요망 02.28
*컨택자 : 김호영
*현재 상황 : 졸업자
*공부 목적+니즈(구체적으로) : 모스 자격증 취득
*전공 유무 및 해본적 있는지 : 없음
*거주지or직장위치+학원까지 이동시간 : 미파악""", 11.2)

exact_page(find("국취제 과정 확대 · 원문 2"), """*평일/주말 스케줄이 어떤지 (가용시간 체크) : 상관 없음
*공부 시작 시기 : 3월
*생각 중인 공부 기간 : 길지 않음
*과정 안내 내용 : 모스과정 안내
*상담 명분 제시+상담 유도 : 국비조회
*같이 방문하실 지인 : 친한 동생
*이외 상담내용 :
이제 막 졸업한 분으로 친한 동생이랑 모스를 배우고자 함 -> 동생은 4학년
국비지원이 되는지 물어보았고 둘 다 가능하다고 안내
모스 자격증을 취득하려하고 특별하게 취득해야할 이유는 없다고 함
따라서 취업지원까지 가능한 국비제도 안내해드리기로 하였음
본인은 시간이 많아서 괜찮은데 동생이랑 시간 맞춰보고 상담일정 잡아보겠다고 함
**수강료 질문에 대한 응대 : X
**시간표 질문에 대한 응대 : 오후, 저녁
*최소 과정: 모스
*최대 과정: 사무패키지
*특이 사항: X
*통화 시간: 6분""", 11.5)

exact_page(find("국취제 과정 확대 · 원문 3"), """컨택결과 : 완납 3,330,000
컨택내용
*상담자 : 김호영
*직업 : 구직자
*나이 : 00년생
*거주지 : 도안동
*과정 관심 계기 : 창업 전 취업
*관심 및 희망 과정 : 모스
*수강 목적 : 취업 시 도움이 될까
*학교&직장 내 특이사항 : X
*평소 가용 가능한 시간 : 저녁 선호 (황예림수강생이랑 같은 시간대 듣기 위해)
*희망 개강 시기 : 3월
*희망 타임 : 저녁
*수강 예산 : 333
*성격 및 성향 : 보통
*결제권자/완납 예정 : 아버님 / 완납
*환불 리스크 : X
*추가 등록 가망 과목/시기 : X
*이외 상담 내용 :""", 11.8)

exact_page(find("국취제 과정 확대 · 원문 4"), """건설안전용품을 판매하는 사업을 하고 싶음
아버님께서는 건설철거 관련해서 사업을 하시는 듯 -> 학생 때 국장 못받음
국취제 받기 위해 주소이전 안내하였고 서울에 있는 언니집으로 옮기는 쪽으로 할 것 같음
건축, 경리, 마케팅 분야로 안내하다가 마케팅쪽으로 유도하여 웹디자인 안내
황예림 수강생이 친구인데 저녁이랑 주말로 수업 다니고 있어서 같은 타임으로 공부하길 희망함
웹디자인으로 완납
**이번 상담 때 나의 반성 및 고칠 점 : X""", 13)


# 26: retain blank source fields; add one page so the text does not have to shrink.
slide_26_1 = find("대상자가 아닐 때 · 원문 1")
exact_page(slide_26_1, """컨택결과 : 상담예정 09.20 ★13:30
컨택일시 : 2022-09-19 12:12:16
컨택자 : 2-1팀 김호영2
컨택내용
*컨택자 : 김호영
*본인확인 및 문의여부체크 : 본인
*학생분 현재 상황 질문 : 휴학생
*배워보시려는 계기 :
*분야 경험 여부 및 교육 여부 :
*현재 거주지 :
*평소 평일/주말 스케줄이 어떤지 (가용시간체크) :
*배우고 싶은 시기 : 11월
*질문한 사유 안내 + 위 내용을 취합 후 정리 :
*알아보는 과정에 대한 안내 :
*상담 명분 제시 + 상담 유도 :
*상담 일정 확정 :
*같이 방문하실 지인 :
*이외 상담내용 :""", 11.8)

old_26_2 = find("대상자가 아닐 때 · 원문 2")
set_title(old_26_2, "대상자가 아닐 때 · 원문 3", "26-3")
slide_26_2 = new_after(slide_26_1, "CHAPTER 3. 상담 스피치", "대상자가 아닐 때 · 원문 2", "26-2")
exact_page(slide_26_2, """24살로 건설관련 학과였는데 현장직이 아닌 사무실내에서 근무하는 것을 해보고 싶어 현재 1년간 휴학중
아직까지 복학할 계획은 없다고 함
웹디자인과 광고디자인 둘 다 관심이 있는 상황이고 국비지원은 휴학생 불가라고 안내 -> 휴학생은 교육청지원 안내
본인도 국비지원이 더 많이 되는거죠라고 물어봐서 그렇다고 하였음
부모님과도 조율 중이라 11월 희망한다고 하였는데 상담왔을 시 굳혀버릴 것""", 13)

exact_page(old_26_2, """휴학생으로 부사관학과인데 그쪽으로 취업하지 않을 거라 크게 관심 없다고 함
친구가 웹프로그래머로 근무하는데 괜찮아 보여서 관심이 생겼고 웹디자인 전망도 좋은 부분 안내
마지막엔 인테리어쪽도 조금 궁금해하며 고민이 다소 있는 듯함
익일까지 고민하기로 했고 국취제 참여하는 쪽으로 고려중
자부담금 설명해줬더니 생각보다 금액이 많이 안나왔다고 하였음

학교 자퇴하지 않고 다니기로 하였고
따라서 국취제는 내년 9월에 신청하기로 함
1개월에 1개씩 교육과정 진행하기로 하였음
김호영2 (7696122)
기타
- COD 기한 : 09.21 (수)
- COD 결제수단 : 카드예정
- 결제권자 : 부모님
- 사유 : 국취제 참여 위해 학교 자퇴 고려중""", 12.5)


# 27: source labels and phrases, including the source's spacing.
exact_page(find("국취제 할인율 · 원문 상담"), """컨택결과 : 완납 3,000,000
컨택일시 : 2023-05-15 20:04:25
컨택자 : 2-1팀 김호영2
컨택내용
*상담자 : 김호영
*직업 : 구직자
*나이 : 26세
*거주지 : 중구 유천동
*과정 관심 계기 : 취업
*관심 및 희망 과정 : 컴활, 기계
*수강 목적 : 취업
*학교&직장 내 특이사항 : X
*평소 가용 가능한 시간 : 프리함
*희망 개강 시기 : 6월
*희망 타임 : 오후 (9시만 아니면 됨)
*수강 예산 : 300
*성격 및 성향 : 보통
*결제권자/완납 예정 : 부모님 / 완납
*환불 리스크 : X
*추가 등록 가망 과목/시기 : X""", 11.5)

exact_page(find("국취제 할인율 · 원문 상담내용"), """*이외 상담 내용 :
어머님과 함께 방문하였음
아들이 졸업하고 취업준비를 안하고 게임만해서 답답해서 전화주심
아들은 일단 컴활 먼저 하고 나중에 봐서 하겠다는 식이었으나 지금 해야한다고 압박을 주었음
어머님도 멘토편이었고 조심스럽게 달래가면서 이야기하였음
결국은 본인도 기계학과를 살려야겠다며 기사 자격증을 준비하는 쪽으로 이야기 함
컴활부터 진행해서 국취제 통해 기계 분야 공부 진행해보기로 하였음""", 13.5)


# Remove any remaining non-cover source pictures and keep the presentation font consistent.
for slide_index, slide in enumerate(prs.slides, 1):
    for shape in list(slide.shapes):
        if slide_index > 1 and shape.shape_type == 13:
            remove(shape)
        if getattr(shape, "has_text_frame", False):
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = FONT
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                for cell in row.cells:
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.name = FONT


prs.save(out)
print(f"saved={out} slides={len(prs.slides)}")
