import sys
from io import BytesIO
from pathlib import Path
sys.path.insert(0, str(Path(__file__).resolve().parents[1] / 'pptdeps'))
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from PIL import Image

path=Path(sys.argv[1]); output=Path(sys.argv[2]) if len(sys.argv)>2 else path; prs=Presentation(path)
FONT='Noto Sans KR'; NAVY=RGBColor(6,70,130); BLUE=RGBColor(10,79,153); DARK=RGBColor(10,53,106)
PALE=RGBColor(238,247,255); PALE2=RGBColor(221,240,255); BORDER=RGBColor(185,217,245); TEXT=RGBColor(30,43,56); WHITE=RGBColor(255,255,255)

def remove(sh): sh._element.getparent().remove(sh._element)
def shape(slide,name): return next((s for s in slide.shapes if s.name==name),None)
def set_box(sh,x,y,w,h): sh.left=Pt(x);sh.top=Pt(y);sh.width=Pt(w);sh.height=Pt(h)
def format_runs(sh,min_size=11.5):
    if getattr(sh,'has_text_frame',False):
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                r.font.name=FONT
                if r.font.size is not None and r.font.size.pt < min_size:r.font.size=Pt(min_size)
    if getattr(sh,'has_table',False):
        for row in sh.table.rows:
            for cell in row.cells:
                for p in cell.text_frame.paragraphs:
                    for r in p.runs:
                        r.font.name=FONT
                        if r.font.size is None or r.font.size.pt < 10.5:r.font.size=Pt(10.5)
def box(slide,x,y,w,h,text,fill=PALE,line=BORDER,size=16,bold=False,color=TEXT,align=PP_ALIGN.LEFT):
    s=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Pt(x),Pt(y),Pt(w),Pt(h));s.fill.solid();s.fill.fore_color.rgb=fill;s.line.color.rgb=line
    tf=s.text_frame;tf.clear();tf.word_wrap=True;tf.margin_left=tf.margin_right=Pt(16);tf.margin_top=tf.margin_bottom=Pt(8);tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    p=tf.paragraphs[0];p.alignment=align;r=p.add_run();r.text=text;r.font.name=FONT;r.font.size=Pt(size);r.font.bold=bold;r.font.color.rgb=color
    return s
def clear_below_header(slide):
    for sh in list(slide.shapes):
        if sh.top.pt > 82: remove(sh)

# Preserve every evidence image before removing it from crowded slides.
evidence=[]
for slide_no in (6,7,10,11,20,23,24,25,26,27):
    sl=prs.slides[slide_no-1]
    for sh in list(sl.shapes):
        if sh.shape_type==13:
            evidence.append((slide_no,sh.name,sh.image.blob))

# Rebuild crowded image slides as large, readable summary slides.
summaries={
6:[('01','내일배움카드 교육을 찾는 경우','국비 수요의 80% 이상은 내일배움카드에 해당합니다. 지원 제도와 과정의 장단점을 비교해 안내합니다.'),('02','국비 언급은 없으나 국취제 조건이 가능한 경우','적용 가능한 국비를 먼저 언급해 타 학원과의 중복을 방지하고 방문 명분을 제시합니다.'),('03','국민취업지원제도로 방문 유도','방문 시 모의 산정을 통해 받을 수 있는 지원금을 확인해드린다고 안내합니다.'),('04','상담 날짜를 당기고 싶은 경우','예산과 참여 절차, 개강 후 참여 제한을 설명해 빠른 상담의 필요성을 안내합니다.')],
7:[('01','내일배움카드 문의자','현재 상태와 카드 발급·사용 여부를 먼저 확인합니다.'),('02','국민취업지원제도 가능 문의자','나이, 재직 여부, 실업급여, 가구원과 소득 조건을 확인합니다.'),('03','제도 비교가 필요한 문의자','지원율만 강조하지 않고 일정·과정 적합성·취업 목표를 함께 비교합니다.'),('04','방문 상담 유도','방문 시 정확한 조회와 모의 산정이 가능하다는 점을 안내합니다.')],
10:[('01','DB 출처에 맞는 스피치','영상, 인테리어, AI, 연령 타깃 등 유입 경로에 맞춰 말합니다.'),('02','현재 상황 파악','학생·구직자·재직자 여부를 우선 확인합니다.'),('03','맞춤 제도 선택','다양한 국비 제도 중 본인에게 맞는 제도가 중요합니다.'),('04','국취제 조건 확인','나이, 재직 여부, 실업급여, 가구원을 미리 확인합니다.'),('05','컨택 클로징','제도를 비교해 최소 자부담 방향을 안내합니다.')],
11:[('01','실업급여와 국비지원','일반과정도 구직활동으로 인정될 수 있음을 안내합니다.'),('02','내일배움카드 경험 확인','고용센터에서 안내받았는지와 기존 사용 여부를 확인합니다.'),('03','자부담 안내','전액지원이 어려울 수 있으므로 자부담 발생 가능성을 설명합니다.'),('04','잔액·납부 내역 확인','기존 카드 잔액과 과거 자부담 납부 내역을 확인합니다.'),('05','추천 방향','지점 국비과정이 있다면 국비와 일반과정의 조합을 추천합니다.')]
}
for slide_no,items in summaries.items():
    sl=prs.slides[slide_no-1];clear_below_header(sl)
    cols=2 if len(items)<=4 else 1
    if cols==2:
        for i,(num,title,body) in enumerate(items):
            x=42+(i%2)*448;y=102+(i//2)*184
            box(sl,x,y,427,165,f'{num}  {title}\n\n{body}',WHITE,BORDER,15,False,TEXT)
    else:
        for i,(num,title,body) in enumerate(items):
            y=98+i*76
            box(sl,42,y,875,66,f'{num}  {title}   |   {body}',PALE if i%2==0 else WHITE,BORDER,13.5,False,TEXT)
    box(sl,122,474,715,28,'관련 원본 이미지는 문서 마지막의 확대 이미지 페이지에서 확인',PALE2,BORDER,11.5,True,DARK,PP_ALIGN.CENTER)

simple={
20:'졸업예정자의 졸업 예정 7개월 전부터 신청 가능\n\n국취제는 취업 대비 목적이 많으므로, 4학년 1학기에는 기초 과정을 먼저 진행하고 2학기 자격증·포트폴리오 시기에 신청해 수당을 지급받는 흐름으로 안내합니다.',
23:'카카오톡을 활용해 상담 전후 자료와 신청 링크를 전달하고, 문의자의 반응과 진행 상태를 이어서 관리합니다.\n\n원본 대화 화면은 마지막 확대 이미지 페이지에서 확인할 수 있습니다.',
24:'기존 수강생의 실제 수급 내역을 근거로 지원금 수령 사례를 설명합니다.\n\n금액과 입금 내역은 마지막 확대 이미지 페이지에서 크게 확인할 수 있습니다.',
25:'단순 컴퓨터활용능력 취업 목표를 정규과정으로 확장합니다.\n\n사무직이라는 넓은 목표를 경리 또는 마케팅 등 구체적인 직무 방향으로 연결합니다.',
26:'휴학생은 국비 대상이 아니더라도 교육청 지원 방식의 할인을 검토할 수 있습니다.\n\n지원 대상이 아니라는 이유로 자퇴를 권하지 않습니다.',
27:'국민취업지원제도를 단순 할인으로만 설명하지 않습니다.\n\n국비지원은 할인이 필수가 아니며, 과정 적합성과 취업 목표를 중심으로 안내합니다.'}
for slide_no,text in simple.items():
    sl=prs.slides[slide_no-1];clear_below_header(sl)
    box(sl,72,125,815,285,text,PALE,BORDER,20,True,DARK,PP_ALIGN.CENTER)
    box(sl,172,438,615,34,'관련 원본 이미지는 문서 마지막의 확대 이미지 페이지에서 확인',PALE2,BORDER,12,True,DARK,PP_ALIGN.CENTER)

# Improve type size throughout core slides without touching title hierarchy.
for idx in range(2,27):
    sl=prs.slides[idx]
    for sh in sl.shapes:
        if 'Chapter' in sh.name or 'Footer' in sh.name: format_runs(sh,8.5)
        elif 'Title' in sh.name and sh.top.pt < 90: format_runs(sh,22)
        else: format_runs(sh,11.5)

# Append each source image on its own full-size slide, matching the 29+ reference style.
blank=prs.slide_layouts[6]
for source_no,name,blob in evidence:
    sl=prs.slides.add_slide(blank)
    # white full-slide background
    bg=sl.background.fill;bg.solid();bg.fore_color.rgb=WHITE
    im=Image.open(BytesIO(blob));iw,ih=im.size
    sw,sh=prs.slide_width,prs.slide_height
    scale=min(sw/iw,sh/ih);w=int(iw*scale);h=int(ih*scale);x=int((sw-w)/2);y=int((sh-h)/2)
    pic=sl.shapes.add_picture(BytesIO(blob),x,y,w,h);pic.name=f'FINAL_IMAGE_APPENDIX_S{source_no}_{name}'

prs.save(output)
print(f'saved={output} slides={len(prs.slides)} appended={len(evidence)}')
