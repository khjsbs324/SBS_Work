import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent / "deps"))
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

src = Path(sys.argv[1])
prs = Presentation(src)
print(f"slides={len(prs.slides)} size={prs.slide_width}x{prs.slide_height}")
for slide_no in (4, 17):
    slide = prs.slides[slide_no - 1]
    print(f"\nSLIDE {slide_no} shapes={len(slide.shapes)}")
    for idx, shape in enumerate(slide.shapes):
        text = ""
        if getattr(shape, "has_text_frame", False):
            text = " / ".join(p.text for p in shape.text_frame.paragraphs)
        print(
            idx,
            shape.shape_type,
            shape.name,
            f"x={shape.left} y={shape.top} w={shape.width} h={shape.height}",
            repr(text[:500]),
        )
