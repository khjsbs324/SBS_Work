import sys
from pathlib import Path
sys.path.insert(0, str(Path(__file__).resolve().parents[1] / 'pptdeps'))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

PPT = Path(sys.argv[1])
prs = Presentation(PPT)
NAVY=RGBColor(6,70,130); BLUE=RGBColor(10,79,153); DARK=RGBColor(10,53,106)
PALE=RGBColor(238,247,255); PALE2=RGBColor(221,240,255); BORDER=RGBColor(185,217,245)
TEXT=RGBColor(30,43,56); WHITE=RGBColor(255,255,255); FONT='Noto Sans KR'

def delete_content(slide):
    keep={'S16_Chapter','S16_Title','S16_PageBadge','S16_HeaderRule',
          'S18_Chapter','S18_Title','S18_PageBadge','S18_HeaderRule',
          'S19_Chapter','S19_Title','S19_PageBadge','S19_HeaderRule',
          'S21_Chapter','S21_Title','S21_PageBadge','S21_HeaderRule'}
    for sh in list(slide.shapes):
        if sh.name not in keep:
            sh._element.getparent().remove(sh._element)

def style_text(shape, size=12, bold=False, color=TEXT, align=PP_ALIGN.LEFT, margin=8):
    tf=shape.text_frame; tf.clear(); tf.word_wrap=True
    tf.margin_left=tf.margin_right=Pt(margin); tf.margin_top=tf.margin_bottom=Pt(4)
    tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    p=tf.paragraphs[0]; p.alignment=align
    r=p.add_run(); r.text=''
    for para in tf.paragraphs:
        for run in para.runs:
            run.font.name=FONT; run.font.size=Pt(size); run.font.bold=bold; run.font.color.rgb=color
    return tf

def box(slide,x,y,w,h,text='',fill=WHITE,line=BORDER,size=12,bold=False,color=TEXT,align=PP_ALIGN.LEFT,radius=True):
    sh=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,Pt(x),Pt(y),Pt(w),Pt(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=fill; sh.line.color.rgb=line; sh.line.width=Pt(1)
    tf=sh.text_frame; tf.clear();tf.word_wrap=True;tf.margin_left=tf.margin_right=Pt(12);tf.margin_top=tf.margin_bottom=Pt(5);tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    p=tf.paragraphs[0];p.alignment=align;r=p.add_run();r.text=text;r.font.name=FONT;r.font.size=Pt(size);r.font.bold=bold;r.font.color.rgb=color
    return sh

def add_label(slide,x,y,w,h,text,fill=NAVY,size=13):
    return box(slide,x,y,w,h,text,fill,fill,size,True,WHITE,PP_ALIGN.CENTER,False)

def add_card(slide,x,y,w,h,num,title,body):
    base=box(slide,x,y,w,h,'',WHITE,BORDER)
    badge=box(slide,x+12,y+12,34,34,num,PALE2,PALE2,11,True,BLUE,PP_ALIGN.CENTER)
    tb=slide.shapes.add_textbox(Pt(x+58),Pt(y+10),Pt(w-70),Pt(h-18))
    tf=tb.text_frame;tf.clear();tf.word_wrap=True;tf.margin_left=tf.margin_right=0;tf.margin_top=tf.margin_bottom=0
    p=tf.paragraphs[0];r=p.add_run();r.text=title;r.font.name=FONT;r.font.size=Pt(13);r.font.bold=True;r.font.color.rgb=DARK
    p=tf.add_paragraph();p.space_before=Pt(5);r=p.add_run();r.text=body;r.font.name=FONT;r.font.size=Pt(11.3);r.font.color.rgb=TEXT

def set_cell(cell,text,fill,size=10,bold=False,color=TEXT,align=PP_ALIGN.CENTER):
    cell.fill.solid();cell.fill.fore_color.rgb=fill
    cell.margin_left=cell.margin_right=Pt(4);cell.margin_top=cell.margin_bottom=Pt(3)
    tf=cell.text_frame;tf.clear();tf.word_wrap=True;tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    p=tf.paragraphs[0];p.alignment=align;r=p.add_run();r.text=str(text);r.font.name=FONT;r.font.size=Pt(size);r.font.bold=bold;r.font.color.rgb=color

def native_table(slide,x,y,w,h,headers,rows,font=10,col_widths=None):
    shape=slide.shapes.add_table(len(rows)+1,len(headers),Pt(x),Pt(y),Pt(w),Pt(h));tbl=shape.table
    if col_widths:
        for i,cw in enumerate(col_widths):tbl.columns[i].width=Pt(cw)
    for c,v in enumerate(headers):set_cell(tbl.cell(0,c),v,NAVY,font,True,WHITE)
    for r,row in enumerate(rows,1):
        fill=PALE if r%2 else WHITE
        for c,v in enumerate(row):set_cell(tbl.cell(r,c),v,fill,font,False,TEXT)
    return shape

# Slide 16: large, readable principle cards and script cards.
s=prs.slides[15];delete_content(s)
add_label(s,42,94,430,34,'내일배움카드 타파 원칙')
add_label(s,490,94,427,34,'기본 컨택 스피치 · 편집 가능한 원문')
principles=[
('원칙 1','좋은 제도라는 점은 인정하되, 문의자의 상황에는 맞지 않는 지원임을 설명합니다.\n문의자의 상황에 맞춰 단점을 구체적으로 안내합니다.'),
('원칙 2','본원의 장점과 멘토의 필요성을 상담에서 충분히 전달하면 내일배움카드에만 머무르지 않도록 설득할 수 있습니다.'),
('원칙 3','지원이 되지 않는 수업은 이유를 명확히 설명합니다.\n특수학과 취업직군, 채용공고가 많지 않은 취업군은 당당하게 안내합니다.')]
for i,(t,b) in enumerate(principles):add_card(s,42,140+i*112,430,101,f'{i+1:02d}',t,b)
scripts=[
('내일배움카드 교육 찾는 경우','국비 수요의 80% 이상은 내일배움카드에 해당합니다. 지원 제도와 과정의 장단점을 비교해 선택하도록 안내합니다.'),
('국비 언급은 없으나 국취제 조건이 가능한 경우','적용 가능한 국비를 먼저 언급해 타 학원과의 중복을 방지하고 방문 명분을 제시합니다.'),
('국민취업지원제도로 방문 유도','방문 시 모의 산정을 통해 받을 수 있는 지원금을 확인해드린다고 안내합니다.'),
('상담 날짜를 당기고 싶은 경우','예산과 참여 절차, 개강 후 참여 제한을 안내해 빠른 상담의 필요성을 설명합니다.')]
for i,(t,b) in enumerate(scripts):add_card(s,490,140+i*84,427,74,f'{i+1:02d}',t,b)

# Slide 18: screenshot comparison converted to readable native table/cards.
s=prs.slides[17];delete_content(s)
add_label(s,42,94,875,34,'국민취업지원제도에 참여할 수 없는 대상')
type1=['취업 중인 자(단, 임금근로자 주 30시간 미만·사업소득자 월소득 250만원 미만은 불완전 취업자로 참여 가능)','근로능력·취업 및 구직의사가 없는 사람','생계급여 수급자','구직급여 수급 중이거나 종료 후 6개월이 지나지 않은 사람','재정지원 일자리사업 참여 중이거나 종료 후 6개월이 지나지 않은 사람','국가·지자체의 구직활동 지원수당 수급 중이거나 종료 후 6개월이 지나지 않은 사람','본인 월평균 소득이 1인 가구 중위소득 60%를 넘는 사람','상급학교 진학·전문자격 취득 목적으로 재학 또는 수강 중인 사람']
type2=['취업 중인 자(단, 임금근로자 주 30시간 미만·사업소득자 월소득 250만원 미만은 불완전 취업자로 참여 가능)','근로능력·취업 및 구직의사가 없는 사람','구직급여 수급자(수급 종료 후 참여 가능)','재정지원 일자리사업 참여 중인 사람(종료 후 참여 가능)','국가·지자체의 구직활동 지원수당 수급 중인 사람(수급 종료 후 참여 가능)','상급학교 진학·전문자격 취득 목적으로 재학 또는 수강 중인 사람']
add_label(s,42,140,424,30,'Ⅰ유형',BLUE,12);add_label(s,493,140,424,30,'Ⅱ유형',BLUE,12)
for x,items in [(42,type1),(493,type2)]:
    text='\n'.join('• '+v for v in items);b=box(s,x,176,424,238,text,PALE,BORDER,9.6,False,TEXT,PP_ALIGN.LEFT)
    b.text_frame.vertical_anchor=MSO_ANCHOR.TOP
box(s,42,426,875,61,'상담 핵심  |  초반 해당 여부 확인 → 지원대상 페이지 확인 → 참여 제한 항목 설명\nⅠ유형 제한 사유: 학생은 즉시 취업이 어렵고, 생계급여·실업급여는 이중 지원에 해당',PALE2,BORDER,11.2,True,DARK,PP_ALIGN.CENTER)

# Slide 19: income screenshot converted to a native table.
s=prs.slides[18];delete_content(s)
add_label(s,42,94,875,34,'26년 기준 중위소득  |  단위: 원')
headers=['구분','1인 가구','2인 가구','3인 가구','4인 가구','5인 가구']
rows=[['중위소득 60%','1,538,543','2,519,575','3,215,422','3,896,843','4,534,031'],['중위소득 100%','2,564,238','4,199,292','5,359,036','6,494,738','7,556,719'],['중위소득 120%','3,077,086','5,039,150','6,430,843','7,793,686','9,068,063']]
native_table(s,62,145,835,142,headers,rows,11.2,[140,139,139,139,139,139])
box(s,62,301,835,63,'가구단위 산정 기준\n신청인 본인, 배우자(사실혼 포함), 1촌 이내 직계혈족(부모·자녀) 중 생계와 주거를 같이 하는 민법상 가족을 포함합니다.',PALE,BORDER,11.3,False,TEXT,PP_ALIGN.LEFT)
items=[('01','가구단위','4인 가구라도 형제·자매는 미포함되어 3인 가구로 산정될 수 있습니다.'),('02','소득 초과 시','주소지 이전(전입신고) 가능 여부를 확인합니다.'),('03','주소지 이전 순서','조부모 → 형제·자매 → 친척 → 친구 순으로 검토합니다.')]
for i,(num,t,b) in enumerate(items):add_card(s,62+i*280,380,265,83,num,t,b)
box(s,202,474,555,29,'정부24를 통한 간편 전입신고 방법을 함께 안내합니다.',PALE2,BORDER,10.8,True,DARK,PP_ALIGN.CENTER)

# Slide 21: estimate screenshots converted into two native tuition tables.
s=prs.slides[20];delete_content(s)
graphics=[['9월 프리미어 프로','1개월','400,000'],['10월 애프터이펙트 기초','1개월','450,000'],['11월 애프터이펙트 활용','1개월','450,000'],['12월 애프터이펙트 심화','1개월','450,000'],['9월 블렌더 기초','1개월','500,000'],['10월 블렌더 활용','1개월','500,000'],['11월 3D 에셋지브러쉬 기초','1개월','600,000'],['12월 3D 에셋지브러쉬 활용','1개월','600,000'],['9월 렌더&라이팅 기초','1개월','600,000'],['10월 렌더&라이팅 활용','1개월','600,000'],['11월 리깅&애니메이션 기초','1개월','600,000'],['12월 리깅&애니메이션 활용','1개월','600,000'],['1월 독립&연출','1개월','600,000'],['12월 포트폴리오(무료과목)','6개월','3,600,000']]
machine=[['4월 캐드 기초','1개월','350,000'],['5월 캐드 활용','1개월','350,000'],['6월 인벤터','1개월','400,000'],['7월 퓨전 360 기초','1개월','400,000'],['8월 퓨전 360 활용','1개월','400,000'],['8월 전산응용기계제도기능사','1개월','450,000'],['9월 사무 기초','1개월','300,000'],['10월 컴퓨터활용능력 1급','1개월','350,000'],['5월 AI 프롬프트 기초(무료)','1개월','600,000'],['6월 AI 프롬프트 활용(무료)','1개월','600,000'],['5월 AI 에이전트 기초(무료)','1개월','600,000'],['6월 AI 에이전트 활용(무료)','1개월','600,000']]
add_label(s,42,94,424,32,'3D그래픽학과 등록금 안내');add_label(s,493,94,424,32,'기계학과 등록금 안내')
native_table(s,42,134,424,270,['과목명','기간','등록금'],graphics,7.8,[252,72,100])
native_table(s,493,134,424,270,['과목명','기간','등록금'],machine,8.0,[252,72,100])
box(s,42,416,424,60,'총 등록금 10,550,000원\n장학·국취제 반영 후 최초 납부 예상: 3,600,000원(국취제) + 2,295,000원(자부담) + 990,000원(온라인)',PALE,BORDER,9.2,True,DARK,PP_ALIGN.CENTER)
box(s,493,416,424,60,'총 등록금 5,400,000원\n장학·국취제 반영 후 납부 예상: 3,000,000원(국취제) + 500,000원(온라인)',PALE2,BORDER,9.5,True,DARK,PP_ALIGN.CENTER)
box(s,242,486,475,24,'국취제를 제외한 실제 자부담금을 중심으로 안내',WHITE,BORDER,10.5,True,DARK,PP_ALIGN.CENTER)

# Force Noto Sans KR throughout edited slides and save in place.
for idx in (15,17,18,20):
    for sh in prs.slides[idx].shapes:
        if getattr(sh,'has_text_frame',False):
            for p in sh.text_frame.paragraphs:
                for r in p.runs:r.font.name=FONT
        if getattr(sh,'has_table',False):
            for row in sh.table.rows:
                for cell in row.cells:
                    for p in cell.text_frame.paragraphs:
                        for r in p.runs:r.font.name=FONT
prs.save(PPT)
print(PPT)
