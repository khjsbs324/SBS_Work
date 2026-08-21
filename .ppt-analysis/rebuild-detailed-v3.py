import json, sys
from copy import deepcopy
from pathlib import Path
sys.path.insert(0,str(Path(__file__).resolve().parents[1]/'pptdeps'))
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN,MSO_ANCHOR
from pptx.dml.color import RGBColor
from PIL import Image

src,out,data_path,assets=map(Path,sys.argv[1:5]);prs=Presentation(src);data=json.loads(data_path.read_text(encoding='utf-8'));orig={i:prs.slides[i-1] for i in range(1,51)}
data3=json.loads((data_path.parent/'slide3-content.json').read_text(encoding='utf-8'))
FONT='Noto Sans KR';NAVY=RGBColor(6,70,130);BLUE=RGBColor(10,79,153);DARK=RGBColor(10,53,106);PALE=RGBColor(238,247,255);PALE2=RGBColor(221,240,255);BORDER=RGBColor(185,217,245);TEXT=RGBColor(30,43,56);WHITE=RGBColor(255,255,255);WARM=RGBColor(255,245,223)

def rm(sh):sh._element.getparent().remove(sh._element)
def clear(sl):
    for sh in list(sl.shapes):rm(sh)
    sl.background.fill.solid();sl.background.fill.fore_color.rgb=WHITE
def text(sl,x,y,w,h,txt,size=14,bold=False,color=TEXT,align=PP_ALIGN.LEFT,name=None,margin=0):
    sh=sl.shapes.add_textbox(Pt(x),Pt(y),Pt(w),Pt(h));sh.name=name or 'Text';tf=sh.text_frame;tf.clear();tf.word_wrap=True;tf.margin_left=tf.margin_right=Pt(margin);tf.margin_top=tf.margin_bottom=Pt(1);tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    for i,line in enumerate(str(txt).split('\n')):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph();p.alignment=align;p.space_before=p.space_after=Pt(0);r=p.add_run();r.text=line;r.font.name=FONT;r.font.size=Pt(size);r.font.bold=bold;r.font.color.rgb=color
    return sh
def box(sl,x,y,w,h,txt='',fill=WHITE,line=BORDER,size=14,bold=False,color=TEXT,align=PP_ALIGN.LEFT,name=None,round=True):
    sh=sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if round else MSO_SHAPE.RECTANGLE,Pt(x),Pt(y),Pt(w),Pt(h));sh.name=name or 'Card';sh.fill.solid();sh.fill.fore_color.rgb=fill;sh.line.color.rgb=line;sh.line.width=Pt(1)
    tf=sh.text_frame;tf.clear();tf.word_wrap=True;tf.margin_left=tf.margin_right=Pt(14);tf.margin_top=tf.margin_bottom=Pt(7);tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    for i,line_txt in enumerate(str(txt).split('\n')):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph();p.alignment=align;p.space_before=p.space_after=Pt(0);r=p.add_run();r.text=line_txt;r.font.name=FONT;r.font.size=Pt(size);r.font.bold=bold;r.font.color.rgb=color
    return sh
def header(sl,chapter,title,badge):
    text(sl,42,9,520,14,chapter,8.5,True,RGBColor(138,166,199),name='Chapter')
    text(sl,42,29,800,42,title,25 if len(title)<18 else 22,True,DARK,name='Title')
    box(sl,878,28,39,28,badge,PALE2,PALE2,10.5,True,BLUE,PP_ALIGN.CENTER,'PageBadge')
    line=sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,Pt(42),Pt(79),Pt(875),Pt(1.4));line.fill.solid();line.fill.fore_color.rgb=BLUE;line.line.fill.background();line.name='HeaderRule'
def rebuild(sl,chapter,title,badge):clear(sl);header(sl,chapter,title,badge)
def move_after(new_slide,anchor):
    ids=prs.slides._sldIdLst
    node=next(n for n in ids if int(n.id)==new_slide.slide_id)
    anchor_node=next(n for n in ids if int(n.id)==anchor.slide_id)
    ids.remove(node);idx=list(ids).index(anchor_node);ids.insert(idx+1,node)
def new_after(anchor,chapter,title,badge):
    sl=prs.slides.add_slide(prs.slide_layouts[6]);move_after(sl,anchor);rebuild(sl,chapter,title,badge);return sl
def add_card(sl,x,y,w,h,num,title_txt,body,size=13.5):
    box(sl,x,y,w,h,'',WHITE,BORDER)
    box(sl,x+12,y+12,36,36,num,PALE2,PALE2,11,True,BLUE,PP_ALIGN.CENTER)
    text(sl,x+60,y+10,w-72,25,title_txt,14,True,DARK)
    text(sl,x+60,y+39,w-72,h-49,body,size,False,TEXT)
def band(sl,x,y,w,h,label,fill=NAVY,size=14):box(sl,x,y,w,h,label,fill,fill,size,True,WHITE,PP_ALIGN.CENTER,round=False)
def table(sl,x,y,w,h,headers,rows,size=11.5,widths=None):
    shp=sl.shapes.add_table(len(rows)+1,len(headers),Pt(x),Pt(y),Pt(w),Pt(h));tb=shp.table
    if widths:
        for i,v in enumerate(widths):tb.columns[i].width=Pt(v)
    for ri,row in enumerate([headers]+rows):
        for ci,val in enumerate(row):
            c=tb.cell(ri,ci);c.fill.solid();c.fill.fore_color.rgb=NAVY if ri==0 else (PALE if ri%2 else WHITE);c.margin_left=c.margin_right=Pt(4);c.margin_top=c.margin_bottom=Pt(3);tf=c.text_frame;tf.clear();tf.vertical_anchor=MSO_ANCHOR.MIDDLE
            p=tf.paragraphs[0];p.alignment=PP_ALIGN.CENTER;r=p.add_run();r.text=str(val);r.font.name=FONT;r.font.size=Pt(size);r.font.bold=(ri==0 or ci==0);r.font.color.rgb=WHITE if ri==0 else TEXT
    return shp
def full_image(sl,img_path,caption=''):
    im=Image.open(img_path);iw,ih=im.size;maxw,maxh=875,385;scale=min(maxw/iw,maxh/ih);w,h=iw*scale,ih*scale;x=42+(875-w)/2;y=94+(385-h)/2
    sl.shapes.add_picture(str(img_path),Pt(x),Pt(y),Pt(w),Pt(h))
    if caption:box(sl,142,488,675,28,caption,PALE2,BORDER,11.5,True,DARK,PP_ALIGN.CENTER)
def two_cards(sl,cards):
    for i,c in enumerate(cards):
        y=102+i*190;box(sl,42,y,875,170,'',WHITE,BORDER);band(sl,42,y,875,34,c['title'],BLUE,14);text(sl,62,y+48,835,108,c['body'],13.3)
def numbered(sl,items,top=98,row_h=72):
    for i,it in enumerate(items):
        if isinstance(it,dict):title_txt,body=it.get('title',f'핵심 {i+1}'),it.get('body','')
        else:title_txt,body=f'핵심 {i+1}',it
        add_card(sl,42,top+i*row_h,875,row_h-9,f'{i+1:02d}',title_txt,body,12.5)

# 3: split training types and card-issuance types so every detail is presentation-sized.
sl=orig[3];rebuild(sl,data3['chapter'],data3['title']+' · 훈련 유형','03-1');band(sl,42,98,875,38,data3['trainingCategory'],NAVY,15)
for i,it in enumerate(data3['trainingItems']):add_card(sl,62,154+i*102,835,90,it['number'],it['title'],it['description'],14)
d3=new_after(sl,data3['chapter'],data3['title']+' · 카드 발급 유형','03-2');band(d3,42,98,875,38,data3['cardCategory'].replace('\n',' '),NAVY,15)
ci=data3['cardItems'];add_card(d3,62,154,835,130,ci[0]['number'],ci[0]['title'],ci[0]['description'],13.5)
add_card(d3,62,302,835,82,'05-1',ci[1]['title']+' · '+ci[1]['types'][0]['label'],ci[1]['types'][0]['description'],13.5)
add_card(d3,62,402,835,88,'05-2',ci[1]['title']+' · '+ci[1]['types'][1]['label'],ci[1]['types'][1]['description'],13)

# 4: eligibility, exclusions, rates, and extra-support lists split across three pages.
sp=data['slides']['4'];sl=orig[4];rebuild(sl,sp['chapter'],sp['title']+' · 지원대상','04-1')
box(sl,92,100,775,55,sp['headline'],PALE2,BORDER,18,True,DARK,PP_ALIGN.CENTER)
box(sl,62,173,835,84,sp['eligibilityTitle']+'\n'+sp['eligibilityLead'],PALE,BORDER,15,True,DARK)
band(sl,62,278,835,34,sp['exclusionTitle']+' · 1',BLUE,14)
for i,v in enumerate(sp['exclusions'][:6]):box(sl,62,321+i*28,835,24,f'{i+1:02d}  {v}',WHITE,BORDER,11.8,False,TEXT)
d4a=new_after(sl,sp['chapter'],sp['title']+' · 지원 제외 대상','04-2');band(d4a,42,98,875,36,sp['exclusionTitle']+' · 2',BLUE,15)
for i,v in enumerate(sp['exclusions'][6:],7):box(d4a,62,151+(i-7)*55,835,46,f'{i:02d}  {v}',PALE if i%2 else WHITE,BORDER,13,False,TEXT)
band(d4a,62,438,835,30,sp['rateTitle'],NAVY,13)
# Full rate details continue on the dedicated extra-support page to avoid small type.
d4b=new_after(d4a,sp['chapter'],sp['title']+' · 훈련비 지원율','04-3');band(d4b,62,101,835,36,sp['rateTitle'],NAVY,15)
for i,v in enumerate(sp['rates']):add_card(d4b,62,153+i*78,835,68,f'{i+1:02d}',f'지원율 {i+1}',v,13.2)
d4c=new_after(d4b,sp['chapter'],sp['title']+' · 추가 지원','04-4');box(d4c,92,101,775,50,sp['extraLead'],PALE2,BORDER,17,True,DARK,PP_ALIGN.CENTER);band(d4c,62,169,835,36,sp['extraTitle'],NAVY,15)
for i,v in enumerate(sp['extras']):box(d4c,62,217+i*32,835,27,f'{i+1:02d}  {v}',PALE if i%2==0 else WHITE,BORDER,11.8,False,TEXT)
text(d4c,62,488,835,15,sp['footer'],8.5,False,RGBColor(123,135,148),PP_ALIGN.CENTER)

# 5: one comparison subject per page, preserving every advantage and disadvantage.
sp=data['slides']['5'];sl=orig[5]
for idx,col in enumerate(sp['columns']):
    page=sl if idx==0 else new_after(sl,sp['chapter'],sp['title']+' · '+col['name'],f'05-{idx+1}')
    if idx==0:rebuild(page,sp['chapter'],sp['title']+' · '+col['name'],'05-1')
    box(page,92,100,775,48,col['tagline'],PALE2,BORDER,17,True,DARK,PP_ALIGN.CENTER)
    band(page,62,163,835,31,'장점',BLUE,14)
    for i,v in enumerate(col['advantages']):
        x=62+(i%2)*420;y=205+(i//2)*72;box(page,x,y,405,62,f'{i+1:02d}  {v}',PALE if i%2==0 else WHITE,BORDER,12.3,False,TEXT)
    band(page,62,354,835,31,'단점',NAVY,14)
    for i,v in enumerate(col['disadvantages']):
        x=62+(i%2)*420;y=396+(i//2)*40;box(page,x,y,405,34,f'• {v}',PALE if i%2==0 else WHITE,BORDER,11.7,False,TEXT)

# 6-7: readable response summaries + immediately adjacent source image page.
common=data['commonContactScript']
for n in (6,7):
    sp=data['slides'][str(n)];sl=orig[n];rebuild(sl,sp['chapter'],sp['title'],f'{n:02d}')
    for i,it in enumerate(common):add_card(sl,42,98+i*92,875,82,f'{i+1:02d}',it['title'],it['body'],12.8)
    box(sl,142,478,675,26,'다음 페이지에서 원본 흐름도를 크게 확인',PALE2,BORDER,11,True,DARK,PP_ALIGN.CENTER)
    d=new_after(sl,sp['chapter'],sp['title']+' · 원본 확대',f'{n:02d}A');full_image(d,assets/sp['referenceImage'],'원본 내용을 확대 배치한 참고 페이지')

# 8-9: four dense speech cards split into two pages.
for n in (8,9):
    sp=data['slides'][str(n)];sl=orig[n];rebuild(sl,sp['chapter'],sp['title']+' · 1',f'{n:02d}-1');two_cards(sl,sp['cards'][:2])
    d=new_after(sl,sp['chapter'],sp['title']+' · 2',f'{n:02d}-2');two_cards(d,sp['cards'][2:])

# 10-11: text summary first, then each evidence image at readable scale.
for n in (10,11):
    sp=data['slides'][str(n)];sl=orig[n];rebuild(sl,sp['chapter'],sp['title'],f'{n:02d}');numbered(sl,sp['items'],98,70)
    if sp.get('highlight'):box(sl,92,463,775,34,sp['highlight'],WARM,BORDER,11.3,True,DARK,PP_ALIGN.CENTER)
    anchor=sl
    for j,img in enumerate(sp['images'],1):
        d=new_after(anchor,sp['chapter'],sp['title']+f' · 원본 {j}',f'{n:02d}{chr(64+j)}');full_image(d,assets/img,'원본 상담 자료 확대');anchor=d

# 12: four cards split for 14pt body readability.
sp=data['slides']['12'];sl=orig[12];rebuild(sl,sp['chapter'],sp['title']+' · 핵심 1', '12-1');two_cards(sl,sp['cards'][:2])
d=new_after(sl,sp['chapter'],sp['title']+' · 핵심 2','12-2');two_cards(d,sp['cards'][2:]);box(d,92,463,775,35,sp['key'],PALE2,BORDER,11.5,True,DARK,PP_ALIGN.CENTER)

# 13: core understanding and three counseling steps separated.
sp=data['slides']['13'];sl=orig[13];rebuild(sl,sp['chapter'],sp['title']+' · 핵심 이해','13-1');band(sl,42,98,875,38,sp['understandingTitle'],NAVY,15)
for i,v in enumerate(sp['understanding']):add_card(sl,62,154+i*102,835,90,f'{i+1:02d}',f'핵심 이해 {i+1}',v,13.5)
d13=new_after(sl,sp['chapter'],sp['title']+' · 상담 단계','13-2')
for i,st in enumerate(sp['steps']):add_card(d13,62,102+i*116,835,104,f'{i+1:02d}',f'{i+1}. {st["title"]}',st['body'],13.2)
box(d13,92,463,775,35,sp['key'],PALE2,BORDER,11.5,True,DARK,PP_ALIGN.CENTER)

# 14: introduction and two detailed response pages.
sp=data['slides']['14'];sl=orig[14];rebuild(sl,sp['chapter'],sp['title']+' · 기본 이해','14-1');text(sl,42,88,875,14,sp['subtitle'],11.3,False,RGBColor(123,135,148));band(sl,42,108,875,34,sp['introTitle']);
for i,v in enumerate(sp['intro']):add_card(sl,62,151+i*90,835,80,f'{i+1:02d}',f'기본 포인트 {i+1}',v,13)
anchor=sl
for i,sc in enumerate(sp['speeches'],1):
    d=new_after(anchor,sp['chapter'],sp['title']+f' · 스피치 {i}',f'14-{i+1}');band(d,42,98,875,36,sc['title'],BLUE if i==1 else NAVY,15)
    box(d,62,151,835,135,'설명\n'+sc['description'],PALE,BORDER,14,False,TEXT)
    box(d,62,303,835,125,'응대 포인트\n'+sc['response'],PALE2,BORDER,14,True,DARK)
    box(d,112,455,735,35,sp['key'],WARM,BORDER,11.3,True,DARK,PP_ALIGN.CENTER);anchor=d

# 15: one concept per page.
sp=data['slides']['15'];sl=orig[15];rebuild(sl,sp['chapter'],sp['title']+' · 스피치 3','15-1');text(sl,42,88,875,14,sp['subtitle'],11.3,False,RGBColor(123,135,148))
box(sl,42,108,875,82,'유형\n'+sp['type'],PALE2,BORDER,14,True,DARK)
for i,v in enumerate(sp['responsePoints']):add_card(sl,42,198+i*86,875,76,f'{i+1:02d}',f'응대 포인트 {i+1}',v,13)
anchor=sl;d=new_after(anchor,sp['chapter'],sp['title']+' · 변수 응대','15-2');band(d,42,98,875,36,sp['variableTitle'],BLUE,15)
for i,q in enumerate(sp['quotes']):box(d,62,158+i*135,835,115,q,PALE if i==0 else PALE2,BORDER,15,False,TEXT)
box(d,62,430,835,63,'활용 팁\n'+'\n'.join('• '+t for t in sp['tips']),WARM,BORDER,11.3,True,DARK)
d2=new_after(d,sp['chapter'],sp['title']+' · 상담 마무리','15-3');numbered(d2,sp['closing'],105,82);box(d2,92,463,775,35,sp['key'],WARM,BORDER,11.5,True,DARK,PP_ALIGN.CENTER)

# 16: principles and scripts separated.
sp=data['slides']['16'];sl=orig[16];rebuild(sl,sp['chapter'],sp['title']+' · 원칙','16-1')
for i,v in enumerate(sp['principles']):add_card(sl,42,105+i*120,875,105,f'{i+1:02d}',f'원칙 {i+1}',v,13.5)
d=new_after(sl,sp['chapter'],sp['title']+' · 기본 컨택 스피치','16-2')
for i,it in enumerate(common):add_card(d,42,100+i*92,875,82,f'{i+1:02d}',it['title'],it['body'],12.8)

# 17: native table with larger type.
sp=data['slides']['17'];sl=orig[17];rebuild(sl,sp['chapter'],sp['title'],'17')
box(sl,42,96,875,38,'✓ '+sp['section1'],PALE2,BORDER,14,True,DARK)
table(sl,62,148,835,212,sp['headers'],sp['rows'],11.5,[209,209,208,209])
box(sl,42,383,875,38,'✓ '+sp['section2'],PALE2,BORDER,14,True,DARK)
box(sl,62,437,835,53,sp['note'],PALE,BORDER,12.5,False,TEXT)

# 18: overview plus two large exclusion panels.
sp=data['slides']['18'];sl=orig[18];rebuild(sl,sp['chapter'],sp['title']+' · 상담 핵심','18-1');numbered(sl,sp['items'],100,78);box(sl,92,460,775,38,sp['highlight'],WARM,BORDER,11.5,True,DARK,PP_ALIGN.CENTER)
type1=['취업 중인 자(단, 임금근로자 주 30시간 미만, 사업소득자 월소득 250만원 미만은 불완전 취업자로 참여 가능)','근로능력, 취업 및 구직의사가 없는 사람','생계급여 수급자','구직급여를 수급 중이거나 수급 종료 후 6개월이 지나지 않은 사람','재정지원 일자리사업에 참여 중이거나 참여 종료 후 6개월이 지나지 않은 사람. 단 일부사업의 경우 참여 종료 후 즉시 참여 가능','국가 또는 지방자치단체가 구직활동에 필요한 비용을 지원하는 수당을 수급 중이거나 수급 종료 후 6개월이 지나지 않은 사람','신청인 본인의 월평균 소득이 1인 가구 기준 중위소득의 60%(1,538,543원)가 넘는 사람','상급학교 진학 및 전문자격증 취득을 목적으로 각종학교에 재학 또는 학원 등에서 수강중인 사람']
type2=['취업 중인 자(단, 임금근로자 주 30시간 미만, 사업소득자 월소득 250만원 미만은 불완전 취업자로 참여 가능)','근로능력, 취업 및 구직의사가 없는 사람','구직급여 수급자(수급 종료 후 참여 가능)','재정지원 일자리사업에 참여 중인 사람(종료후 참여가능)','국가 또는 지방자치단체가 구직활동에 필요한 비용을 지원하는 수당을 수급 중인 사람(수급 종료 후 참여 가능)','상급학교 진학 및 전문자격증 취득을 목적으로 각종학교에 재학 또는 학원 등에서 수강중인 사람']
anchor=sl;page=2
for lab,items in (('Ⅰ유형',type1),('Ⅱ유형',type2)):
    for start in range(0,len(items),4):
        batch=items[start:start+4];d=new_after(anchor,sp['chapter'],sp['title']+f' · {lab}',f'18-{page}');band(d,42,98,875,38,f'{lab} 참여 제한 대상',BLUE,15)
        for j,v in enumerate(batch):add_card(d,62,154+j*82,835,70,f'{start+j+1:02d}',f'{lab} 제한 항목 {start+j+1}',v,12.8)
        anchor=d;page+=1
d=new_after(anchor,sp['chapter'],sp['title']+' · 원본 기준표','18-A');full_image(d,assets/sp['image'],'국민취업지원제도 참여 제한 대상 원본 확대')

# 19: table and counseling guide separated.
sp=data['slides']['19'];sl=orig[19];rebuild(sl,sp['chapter'],sp['title']+' · 기준표','19-1');band(sl,42,98,875,36,'26년 기준 중위소득  |  단위: 원')
headers=['구분','1인 가구','2인 가구','3인 가구','4인 가구','5인 가구'];rows=[['중위소득 60%','1,538,543','2,519,575','3,215,422','3,896,843','4,534,031'],['중위소득 100%','2,564,238','4,199,292','5,359,036','6,494,738','7,556,719'],['중위소득 120%','3,077,086','5,039,150','6,430,843','7,793,686','9,068,063']]
table(sl,62,154,835,170,headers,rows,12,[140,139,139,139,139,139]);box(sl,62,346,835,112,'가구단위 산정 기준\n신청인 본인, 배우자(사실혼 포함), 1촌 이내 직계혈족(부모·자녀) 중 생계와 주거를 같이 하는 민법상 가족을 포함합니다.',PALE,BORDER,13,False,TEXT)
d=new_after(sl,sp['chapter'],sp['title']+' · 상담 가이드','19-2');numbered(d,sp['items'],102,80);box(d,112,445,735,42,sp['highlight'],WARM,BORDER,12,True,DARK,PP_ALIGN.CENTER)
anchor=d
for j,img in enumerate(sp['images'],1):
    di=new_after(anchor,sp['chapter'],sp['title']+f' · 원본 {j}',f'19-A{j}');full_image(di,assets/img,'가구단위·소득 기준 원본 확대');anchor=di

# 20: readable summary followed by the source notice image.
sp=data['slides']['20'];sl=orig[20];rebuild(sl,sp['chapter'],sp['title'],'20');box(sl,92,115,775,295,sp['note'],PALE,BORDER,18,True,DARK,PP_ALIGN.CENTER);box(sl,172,447,615,32,'다음 페이지에서 원본 안내문을 확대 확인',PALE2,BORDER,11.5,True,DARK,PP_ALIGN.CENTER)
d=new_after(sl,sp['chapter'],sp['title']+' · 원본 확대','20A');full_image(d,assets/sp['image'],'졸업예정자 신청 가능 시점 원본')

# 21: each tuition table gets a full page.
sp=data['slides']['21'];graphics=[['9월 프리미어 프로','1개월','400,000'],['10월 애프터이펙트 기초','1개월','450,000'],['11월 애프터이펙트 활용','1개월','450,000'],['12월 애프터이펙트 심화','1개월','450,000'],['9월 블렌더 기초','1개월','500,000'],['10월 블렌더 활용','1개월','500,000'],['11월 3D 에셋지브러쉬 기초','1개월','600,000'],['12월 3D 에셋지브러쉬 활용','1개월','600,000'],['9월 렌더&라이팅 기초','1개월','600,000'],['10월 렌더&라이팅 활용','1개월','600,000'],['11월 리깅&애니메이션 기초','1개월','600,000'],['12월 리깅&애니메이션 활용','1개월','600,000'],['1월 독립&연출','1개월','600,000'],['12월 포트폴리오(무료과목)','6개월','3,600,000']]
machine=[['4월 캐드 기초','1개월','350,000'],['5월 캐드 활용','1개월','350,000'],['6월 인벤터','1개월','400,000'],['7월 퓨전 360 기초','1개월','400,000'],['8월 퓨전 360 활용','1개월','400,000'],['8월 전산응용기계제도기능사','1개월','450,000'],['9월 사무 기초','1개월','300,000'],['10월 컴퓨터활용능력 1급','1개월','350,000'],['5월 AI 프롬프트 기초(무료)','1개월','600,000'],['6월 AI 프롬프트 활용(무료)','1개월','600,000'],['5월 AI 에이전트 기초(무료)','1개월','600,000'],['6월 AI 에이전트 활용(무료)','1개월','600,000']]
sl=orig[21];rebuild(sl,sp['chapter'],sp['title']+' · 3D그래픽학과 1','21-1');table(sl,42,110,875,300,['과목명','기간','등록금'],graphics[:7],12,[600,100,175]);box(sl,112,445,735,38,'3D그래픽학과 등록금 세부 내역 · 1',PALE2,BORDER,13,True,DARK,PP_ALIGN.CENTER)
d=new_after(sl,sp['chapter'],sp['title']+' · 3D그래픽학과 2','21-2');table(d,42,110,875,300,['과목명','기간','등록금'],graphics[7:],12,[600,100,175]);box(d,112,445,735,38,'총 등록금 10,550,000원  |  국취제 제외 실 자부담금 확인',PALE2,BORDER,13,True,DARK,PP_ALIGN.CENTER)
d2=new_after(d,sp['chapter'],sp['title']+' · 기계학과 1','21-3');table(d2,42,110,875,300,['과목명','기간','등록금'],machine[:6],12,[600,100,175]);box(d2,112,445,735,38,'기계학과 등록금 세부 내역 · 1',PALE2,BORDER,13,True,DARK,PP_ALIGN.CENTER)
d3=new_after(d2,sp['chapter'],sp['title']+' · 기계학과 2','21-4');table(d3,42,110,875,300,['과목명','기간','등록금'],machine[6:],12,[600,100,175]);box(d3,112,445,735,38,'총 등록금 5,400,000원  |  국취제 제외 실 자부담금 확인',PALE2,BORDER,13,True,DARK,PP_ALIGN.CENTER)
anchor=d3
for j,img in enumerate(sp['images'],1):
    di=new_after(anchor,sp['chapter'],sp['title']+f' · 원본 견적 {j}',f'21-A{j}');full_image(di,assets/img,sp['callout']);anchor=di

# 22: split payment examples.
sp=data['slides']['22'];sl=orig[22];rebuild(sl,sp['chapter'],sp['title']+' · 자부담 없음','22-1');table(sl,42,112,875,230,sp['headers1'],sp['rows1'],11,[140,105,105,105,105,105,105,105]);box(sl,172,382,615,70,sp['callout1']+'\n'+sp['callout2'],PALE2,BORDER,16,True,DARK,PP_ALIGN.CENTER)
d=new_after(sl,sp['chapter'],sp['title']+' · 자부담 150만원','22-2');table(d,42,112,875,260,sp['headers2'],sp['rows2'],10.8,[140,105,105,105,105,105,105,105]);box(d,172,405,615,70,sp['callout1']+'\n월 납부 금액 250,000원',PALE2,BORDER,16,True,DARK,PP_ALIGN.CENTER)

# 23-27: concise context page, then each image immediately after at maximum readable size.
context={23:'상담 시 카톡 활용\n\n카카오톡을 활용해 상담 전후 자료와 신청 링크를 전달하고, 문의자의 반응과 진행 상태를 이어서 관리합니다.',24:'기존 수강생 수급내역\n\n기존 수강생의 실제 수급 내역을 근거로 지원금 수령 사례를 설명합니다.',25:'단순 컴활 취업 -> 정규 확장\n사무직? -> 경리 or 마케팅\n\n넓은 취업 목표를 구체적인 직무와 정규과정으로 연결합니다.',26:'휴학생은 국비 불가 · 교육청 지원(할인) 가능\n\n자퇴를 권하지 말 것!\n지원 대상이 아니더라도 다른 지원 방식을 검토합니다.',27:'국취제 = 할인\n\n국비지원은 할인이 필수가 아니다!\n과정 적합성과 취업 목표를 중심으로 안내합니다.'}
for n in range(23,28):
    sp=data['slides'][str(n)];sl=orig[n];rebuild(sl,sp['chapter'],sp['title'],str(n));box(sl,82,128,795,260,context[n],PALE,BORDER,20,True,DARK,PP_ALIGN.CENTER);box(sl,172,430,615,38,'이어지는 페이지에서 관련 원본 자료를 한 장씩 확대 확인',PALE2,BORDER,12,True,DARK,PP_ALIGN.CENTER)
    anchor=sl
    for j,img in enumerate(sp['images'],1):
        d=new_after(anchor,sp['chapter'],sp['title']+f' · 원본 {j}',f'{n}-{j}');full_image(d,assets/img,'원본 자료 확대');anchor=d

# Untouched 3-5 and 13: keep design, enforce font consistency without forcing overflow.
prs.save(out);print(f'saved={out} slides={len(prs.slides)}')
