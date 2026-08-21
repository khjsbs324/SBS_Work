import sys
from pathlib import Path
sys.path.insert(0,str(Path(__file__).resolve().parents[1]/'pptdeps'))
from pptx import Presentation
prs=Presentation(sys.argv[1]);
for n,names in ((2,{'TextBox 2','TextBox 50'}),(26,{'Text','Card'})):
 print('SLIDE',n)
 for s in prs.slides[n-1].shapes:
  if s.name in names:print(s.name,round(s.left.pt,1),round(s.top.pt,1),round(s.width.pt,1),round(s.height.pt,1),repr(getattr(s,'text','')))
