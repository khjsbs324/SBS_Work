import sys
from pathlib import Path
sys.path.insert(0,str(Path(__file__).resolve().parents[1]/'pptdeps'))
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN,MSO_ANCHOR
from pptx.dml.color import RGBColor

src,out=map(Path,sys.argv[1:3]);prs=Presentation(src)
FONT='Noto Sans KR';NAVY=RGBColor(6,70,130);BLUE=RGBColor(10,79,153);DARK=RGBColor(10,53,106);PALE=RGBColor(238,247,255);PALE2=RGBColor(221,240,255);BORDER=RGBColor(185,217,245);TEXT=RGBColor(30,43,56);WHITE=RGBColor(255,255,255);WARM=RGBColor(255,245,223)

def title_of(sl):
 vals=[]
 for s in sl.shapes:
  if getattr(s,'has_text_frame',False) and s.text.strip() and s.top.pt<90 and (s.name=='Title' or 'Title' in s.name):vals.append(s.text.strip())
 return max(vals,key=len) if vals else ''
def find(title):return next(s for s in prs.slides if title_of(s)==title)
def rm(s):s._element.getparent().remove(s._element)
def clear(sl):
 for s in list(sl.shapes):rm(s)
 sl.background.fill.solid();sl.background.fill.fore_color.rgb=WHITE
def text(sl,x,y,w,h,txt,size=12,bold=False,color=TEXT,align=PP_ALIGN.LEFT):
 s=sl.shapes.add_textbox(Pt(x),Pt(y),Pt(w),Pt(h));tf=s.text_frame;tf.clear();tf.word_wrap=True;tf.margin_left=tf.margin_right=Pt(1);tf.margin_top=tf.margin_bottom=Pt(1);tf.vertical_anchor=MSO_ANCHOR.TOP
 for i,line in enumerate(str(txt).split('\n')):
  p=tf.paragraphs[0] if i==0 else tf.add_paragraph();p.alignment=align;p.space_after=Pt(1);r=p.add_run();r.text=line;r.font.name=FONT;r.font.size=Pt(size);r.font.bold=bold;r.font.color.rgb=color
 return s
def box(sl,x,y,w,h,txt='',fill=WHITE,line=BORDER,size=12,bold=False,color=TEXT,align=PP_ALIGN.LEFT):
 s=sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Pt(x),Pt(y),Pt(w),Pt(h));s.fill.solid();s.fill.fore_color.rgb=fill;s.line.color.rgb=line;s.line.width=Pt(1);tf=s.text_frame;tf.clear();tf.word_wrap=True;tf.margin_left=tf.margin_right=Pt(14);tf.margin_top=tf.margin_bottom=Pt(7);tf.vertical_anchor=MSO_ANCHOR.TOP
 for i,line in enumerate(str(txt).split('\n')):
  p=tf.paragraphs[0] if i==0 else tf.add_paragraph();p.alignment=align;p.space_after=Pt(1);r=p.add_run();r.text=line;r.font.name=FONT;r.font.size=Pt(size);r.font.bold=bold;r.font.color.rgb=color
 return s
def header(sl,chapter,title,badge):
 text(sl,42,9,520,14,chapter,8.5,True,RGBColor(138,166,199));text(sl,42,29,800,42,title,22 if len(title)>18 else 25,True,DARK)
 b=box(sl,878,28,39,28,badge,PALE2,PALE2,10.5,True,BLUE,PP_ALIGN.CENTER);b.text_frame.vertical_anchor=MSO_ANCHOR.MIDDLE
 l=sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,Pt(42),Pt(79),Pt(875),Pt(1.4));l.fill.solid();l.fill.fore_color.rgb=BLUE;l.line.fill.background()
def rebuild(sl,chapter,title,badge):clear(sl);header(sl,chapter,title,badge)
def move_after(new,anchor):
 ids=prs.slides._sldIdLst;node=next(n for n in ids if int(n.id)==new.slide_id);a=next(n for n in ids if int(n.id)==anchor.slide_id);ids.remove(node);ids.insert(list(ids).index(a)+1,node)
def new_after(anchor,chapter,title,badge):
 s=prs.slides.add_slide(prs.slide_layouts[6]);move_after(s,anchor);rebuild(s,chapter,title,badge);return s
def exact_page(sl,content,size=12):box(sl,42,98,875,398,content,WHITE,BORDER,size,False,TEXT)
def native_table(sl,x,y,w,h,headers,rows,size=11.5,widths=None):
 sh=sl.shapes.add_table(len(rows)+1,len(headers),Pt(x),Pt(y),Pt(w),Pt(h));tb=sh.table
 if widths:
  for i,v in enumerate(widths):tb.columns[i].width=Pt(v)
 for ri,row in enumerate([headers]+rows):
  for ci,v in enumerate(row):
   c=tb.cell(ri,ci);c.fill.solid();c.fill.fore_color.rgb=NAVY if ri==0 else(PALE if ri%2 else WHITE);c.margin_left=c.margin_right=Pt(4);c.margin_top=c.margin_bottom=Pt(3);tf=c.text_frame;tf.clear();tf.word_wrap=True;tf.vertical_anchor=MSO_ANCHOR.MIDDLE
   p=tf.paragraphs[0];p.alignment=PP_ALIGN.CENTER;r=p.add_run();r.text=str(v);r.font.name=FONT;r.font.size=Pt(size);r.font.bold=(ri==0 or ci==0);r.font.color.rgb=WHITE if ri==0 else TEXT
 return sh

# s06: exact flowchart wording, split only at original branches.
sl=find('국비 문의의 흐름');ch='CHAPTER 2. 컨택 스피치';rebuild(sl,ch,'국비 문의의 흐름 · 원문 1','06-1')
exact_page(sl,'''“국비되나요?”\n카드발급 받으셨어요?\n\n사용해보셨을까요? → 네\n사용해봤던 사람들은 국비의 단점 (피드백안되는점, 일을병행할수없는점, 수업기간이 긴점) 을 잘 알고있기때문에, 이 부분만 살짝만 건들어준다면, 국비타파가 쉽다\n\nex) 국비수업을 들어봐서 아시겠지만, 수업시간이 고정이 되어있어서 지금상황에서는 활용이 안되실거에요. 그러면 본원자체에 국비카드 소지자 분들은 저희가 따로 20%~40% 정도 할인적용해서 자부담금내시고 수업 들으실수있으세요''',12.5)
d=new_after(sl,ch,'국비 문의의 흐름 · 원문 2','06-2');exact_page(d,'''사용해보셨을까요? → 아니요\n앞서 말했듯이 카드는 발급을 받았지만 활용을 하지 못하는 경우가 많다. 이럴경우 문의자의 수업 목적, 시간대, 기간을 파악하여 반대로 국비카드를 지금도 역시나 사용하지 못하신다라는 인식을 심어줄것!\n\n문의자 상황 : 알바를 하면서 국비카드 사용되나요? / 직장을 다니고있는데 주말반이나, 저녁반으로 가능한가요?\nex) 국비수업은 9시 ~ 16시까지 수업해야 활용가능하신데 그 시간 수업 들어오실수있으세요?\n=> 직장인들 참여불가\n\n국비수업은 취업위주 수업이기때문에 저녁반이나, 주말반은 수업이 없습니다 => 직장인들 참여불가''',12)
d2=new_after(d,ch,'국비 문의의 흐름 · 원문 3','06-3');exact_page(d2,'''발급신청은 하셨어요? → 네\n내배카로 신청하셨을까요, 국민취업지원제도로 신청하셨을까요?\n먼저는 국비종류 파악이 먼저! 내배카로 신청했을경우 90%이상 그럴경우, 국취제가 좋고, 활용하실수있는 범위가 넓다고 하여서 국취제로 신청하게끔 진행\n☆ 국취제의 장점 어필\n\nex) 내배카 활용하실려면 지금 일 그만두시고 오전부터 낮까지 6개월동안 수업 들어가야하는데 다른일을 병행을 할수없기때문에, 국취제를 활용해서 이직or취업준비하시면서 간단한 알바는 다른 취업준비를 하실수있기때문에 국취제로 신청하면됩니다~\n\n발급신청은 하셨어요? → 아니요\n국비의 종류가 많기때문에, 이야기를 통해서 본인에게 맞는 국비 나라지원을 파악해줄께요\n\n★ 제일쉬운 컨택 및 상담, 등록의 상황''',11.7)

# s07 exact text.
sl=find('국비 문의자 유형');rebuild(sl,ch,'국비 문의자 유형 · 원문','07')
exact_page(sl,'''국비문의자 유형\n\n내일배움카드 1회이상 사용자 15%\n내일배움카드 소지자 40%\n내일배움카드 미소지자 45%\n\n“국비” 겁먹지말고, 자신있게 컨택하자\n\n“국비지원제도, 국비가능한가요?”\n우리가 어떠한 물건을 살때 “혹시 현금으로 하면 할인되나요?”라는 질문과 동일!\n만약, 현금으로 할인이 될 경우에는 안살것인가?\nNO! 구매목적이 있었다면, 현금과 카드는 의미는 없다\n그러나, 소비자 입장에서는 좀더 할인되는 방법을 찾고 싶을뿐\n없다면 그대로 구매한다\n\n실제로 문의준 DB를 컨택을 해보면, 내배카 or 국취제를 정확히 아는 문의자는 없다.\n배움에 목적이 있다면, 국비는 크게 문제될건 없다.''',12.5)

# s10 exact DB table and consultation log.
sl=find('페이스북 DB 컨택 스피치');rebuild(sl,ch,'페이스북 DB 컨택 스피치 · 원문 DB','10-1')
box(sl,42,100,875,178,'''지점 : 대전\n이름 : 함승범(4325163)\n상담결과 : 완납 6,030,000\n핸드폰 : 010-4802-1595\n문의과목 : 드로잉\nDB등록일자 기준 : 2025-08-01 오후 5:00:13\n담당자 : [대전 추천2-1]김호영2\n학과 : 아트윅\n출처 : 디비추천(함승보 010-8543-1595)''',PALE,BORDER,11.8,False,TEXT)
box(sl,42,296,875,178,'''지점 : 대전\n이름 : 함승보(4324482)\n상담결과 : 완납 6,930,000\n핸드폰 : 010-8543-1595\n문의과목 : 영상편집\nDB등록일자 기준 : 2025-08-01 오전 3:17:45\n담당자 : [대전페북2-1]김호영2\n학과 : 모션/영상\n출처 : [페이스Meta]_(대전) 국비 영상편집''',WHITE,BORDER,11.8,False,TEXT)
d=new_after(sl,ch,'페이스북 DB 컨택 스피치 · 원문 상담','10-2');exact_page(d,'''상담예정 08.01 ★16:00\n*컨택자 : 김호영\n*본인확인 및 문의여부 체크 : 본인\n*배워보시려는 계기 및 수강 희망 시기 : 취미 및 취업\n*분야 경험 및 수강 경험 여부 : 프리미어 다뤄봄\n*현재 거주지 및 평소 스케줄 어떤지 (가용시간체크) : 태평동\n*상담 명분 제시 및 상담 일정 확정 : 국비체크\n*이외 상담내용 :\n미디어콘텐츠학과 1학년 다니다가 자퇴했다고 함 → 적성은 맞으나 학교가 맞지 않음\n프리미어는 조금 다룰 수 있고 나머진 해보기만 했다고 하였음\n군대는 가지 않았지만 아직 예정은 없다고 하였음 → 취업, 취미 둘 다 생각중이라고 함\n2D 영상편집 외에도 다른 과정들이 많다고 안내해드렸고 국비 과정 설명해드리기로 하였음\n2시에 둔산동 헬스 레슨이 있어서 온다고 했으며 오전 수업은 선호하지 않음\n9월부터 수강 희망하고 오전 수업 희망하지 않아 영상계좌제 비교하면서 설명할 예정\n**수강료 / 시간표에 대한 나의 응대 : X\n*같이 방문하실 지인 / 기타 특이사항 : X\n*통화시간 : 6분''',11.3)

# s11 exact consultation log and chat.
sl=find('실업급여 수령자 컨택 스피치');rebuild(sl,ch,'실업급여 수령자 컨택 스피치 · 원문 상담','11-1')
exact_page(sl,'''*컨택자 : 김호영\n*본인확인 및 문의여부 체크 : 본인\n*배워보시려는 계기 및 수강 희망 시기 : 취업\n*분야 경험 및 수강 경험 여부 : X\n*현재 거주지 및 평소 스케줄 어떤지 (가용시간체크) : 대전\n*상담 명분 제시 및 상담 일정 확정 : 익일 2시\n*이외 상담내용 :\n유선으로 연결되었고 영상이 아닌 웹디자인 국비 궁금하다고 함\n1년 전에 뷰티과정 1년 국비로 수강했고 추가로 금액도 지불했다고 함\n전에 근무하면서 가볍게 디자인을 경험했던 일이 있었는데 잘 맞는 것 같아서 해보려고 함\n실업급여 수령 중인데 이제 2회차를 받고 있다고 함 → 매월 5일 기준이고 6개월 받는다고 함\n실업급여랑 국비랑 병행 안되는지 여쭤봐서 실업급여 받기 위해서는 구직활동 필요하니까 수업 듣는거 좋다고 안내\n내배카 잔액이랑 체크해드리기 위해 익일 상담 도와드리기로 함\n**수강료 / 시간표에 대한 나의 응대 : X\n*같이 방문하실 지인 / 기타 특이사항 : X\n*통화시간 : 7분''',11.5)
d=new_after(sl,ch,'실업급여 수령자 컨택 스피치 · 원문 대화','11-2');exact_page(d,'''★박지유(오피스)\n멘토님 저 실업급여 제출서류가 필요한데 수강 신청서랑 mos, 지금 듣고있는 컴활2급 출석부 좀 보내주세요~!\n오후 1:48\n\n몇 일부터 몇 일까지 필요하실까요?\n오후 1:48\n\n★박지유(오피스)\nMos랑 컴활 처음 수강날부터 내일 마지막날까지요!\n오후 1:49''',14)

# s19 exact additional source sentences on the guide page.
sl=find('가구단위·소득 기준 안내 · 상담 가이드');rebuild(sl,'CHAPTER 3. 상담 스피치','가구단위·소득 기준 안내 · 원문','19-2')
exact_page(sl,'''• 가구단위: 주민등록표에 등재된 신청인 본인, 신청인의 배우자, 신청인의 1촌 이내 직계 혈족(부모, 자녀)으로 한정\n\n★권준혁(실내건축)\n멘토님 취업 그쪽에서 전화 왔는데 할머니 가구원 추가하면 10만원 더 준다고 하는데 할머니 서명 받으면 될까요?''',15)

# s20 exact notice.
sl=find('졸업예정자 상담');rebuild(sl,'CHAPTER 3. 상담 스피치','졸업예정자 상담 · 원문','20')
exact_page(sl,'''국취 신청기준\n\n@국민취업지원제도 참여제한: 주30시간이상 고용보험가입자/실업급여수급자(2유형은 수급종료후다음날부터 참여가능)/기초생계수급자(본인, 가족포함)(단, 교육, 의료, 주거급여 수급자는 가능), 26년 졸업예정자중 2월졸업자(25.07.01부터), 26년 8월졸업예정자(25.11.01부터)/진학예정자/인문계고 졸업예정자는 졸업년도1.1부터(특성화고는 3학년 당해연도7.1부터/1유형 신청자 본인 소득이 1,435,208(1인 중위소득 60%이상) 이상시 1유형신청불가/사업자등록증 보유자는 월 250만원 미만 소득 또는 월 1,250만원 미만 매출 발생\n\n※ 1유형 소득감액기준: 2025년 기준금액 1,435,208원 (구촉 월지급액 50~70만원인경우 935,208이상 소득발생시 부지급/구촉 월지급액 80만원인경우 80만원이상 소득발생시 부지급)''',11.8)

# s21 exact curricula and scholarship figures.
g1=find('견적서 활용 · 3D그래픽학과 1');rebuild(g1,'CHAPTER 3. 상담 스피치','견적서 활용 · 3D그래픽학과 1','21-1')
native_table(g1,42,105,875,330,['월','과목명','COURSE','등록금'],[['9월','프리미어 프로','1개월','400,000'],['10월','에프터이펙트 기초','1개월','450,000'],['11월','에프터이펙트 활용','1개월','450,000'],['12월','에프터이펙트 심화','1개월','450,000'],['9월','블렌더 기초','1개월','500,000'],['10월','블렌더 활용','1개월','500,000'],['11월','3D 에셋&브러쉬 기초','1개월','600,000']],12,[90,500,120,165])
g2=find('견적서 활용 · 3D그래픽학과 2');rebuild(g2,'CHAPTER 3. 상담 스피치','견적서 활용 · 3D그래픽학과 2','21-2')
native_table(g2,42,105,875,330,['월','과목명','COURSE','등록금'],[['12월','3D 에셋&브러쉬 활용','1개월','600,000'],['9월','렌더&라이팅 기초','1개월','600,000'],['10월','렌더&라이팅 활용','1개월','600,000'],['11월','리깅&애니메이션 기초','1개월','600,000'],['12월','리깅&애니메이션 활용','1개월','600,000'],['1월','룩덥&언리얼','1개월','600,000'],['12월','포트폴리오 (무료과목)','6개월','3,600,000']],12,[90,500,120,165]);box(g2,142,458,675,30,'총 등록금 10,550,000',PALE2,BORDER,13,True,DARK,PP_ALIGN.CENTER)
m1=find('견적서 활용 · 기계학과 1');rebuild(m1,'CHAPTER 3. 상담 스피치','견적서 활용 · 기계학과 1','21-3')
native_table(m1,42,105,875,330,['월','과목명','COURSE','등록금'],[['4월','캐드 기초','1개월','350,000'],['5월','캐드 활용','1개월','350,000'],['6월','인벤터','1개월','400,000'],['7월','퓨전 360 기초','1개월','400,000'],['8월','퓨전 360 활용','1개월','400,000'],['8월','전산응용기계제도기능사','1개월','450,000'],['9월','사무 기초','1개월','300,000']],12,[90,500,120,165])
m2=find('견적서 활용 · 기계학과 2');rebuild(m2,'CHAPTER 3. 상담 스피치','견적서 활용 · 기계학과 2','21-4')
native_table(m2,42,105,875,330,['월','과목명','COURSE','등록금'],[['10월','컴퓨터활용능력 1급','1개월','350,000'],['5월','AI 프롬프트 기초 (무료과목)','1개월','600,000'],['6월','AI 프롬프트 활용 (무료과목)','1개월','600,000'],['5월','AI 에이전트 기초 (무료과목)','1개월','600,000'],['6월','AI 에이전트 활용 (무료과목)','1개월','600,000'],['','재수강','1개월','-'],['','재수강','1개월','-']],12,[90,500,120,165]);box(m2,142,458,675,30,'총 등록금 5,400,000',PALE2,BORDER,13,True,DARK,PP_ALIGN.CENTER)
sg=new_after(m2,'CHAPTER 3. 상담 스피치','견적서 활용 · 3D그래픽 장학지원','21-5');native_table(sg,42,112,875,260,['세부장학지원 내역','지원금액','등록금'],[['AIX제이 (~08.12)','3,300,000','7,250,000'],['온라인결합','1,000,000','6,250,000'],['당일등록','355,000','5,895,000'],['국민취업지원제도 1유형','3,600,000','2,295,000']],12,[500,175,200]);box(sg,92,410,775,62,'최초 납부해야할 등록금\n3,600,000(국취제) + 2,295,000(자부담) + 990,000(온라인)',PALE2,BORDER,14,True,DARK,PP_ALIGN.CENTER)
sm=new_after(sg,'CHAPTER 3. 상담 스피치','견적서 활용 · 기계학과 장학지원','21-6');native_table(sm,42,112,875,230,['세부장학지원 내역','지원금액','등록금'],[['새학기 프로모션 (~04.13)','1,800,000','3,600,000'],['온라인결합','600,000','3,000,000'],['국민취업지원제도 1유형','3,600,000','-']],12,[500,175,200]);box(sm,92,390,775,62,'총 납부해야할 등록금\n3,000,000(국취제) + 500,000(온라인)',PALE2,BORDER,14,True,DARK,PP_ALIGN.CENTER)

# s23 exact conversations, one page per original capture.
sl=find('카카오톡 활용');rebuild(sl,'CHAPTER 3. 상담 스피치','카카오톡 활용 · 원문 1','23-1')
exact_page(sl,'''https://blog.naver.com/pj3503/222485497313\n주민등록 주소지 이전변경\n주민등록 주소지 이전 변경 방법 인터넷으로 할수있습니다.\nblog.naver.com\n오후 3:28\n\n국민취업지원...신청 방법.pdf\n올리기일 ~2022. 8. 9. / 용량 387 KB\n오후 3:29\n\n중원님 주소지 옮기고 국민취업지원제도 1유형으로 신청하시면 되세요~ 신청하시다가 헷갈리는 부분 있으시면 연락주세요\n오후 3:43\n\n★청중원(고팅)\n주소지 이전 신청은 했는데\n오후 4:22\n\n★청중원(고팅)\n세대주가 확인을 해야하는데 이거 인증을 하려면 누나가 시간이 아직 안되서 오늘안에 국비지원신청은 못할것 같은데 선생님은 내일 가능한지 천천히 신청해도 괜찮을가요?\n오후 4:23''',11.5)
d=new_after(sl,'CHAPTER 3. 상담 스피치','카카오톡 활용 · 원문 2','23-2');exact_page(d,'''★권준혁(실내건축)\n통화 끝났고 등록 할게요\n오후 2:58\n\nhttps://blog.naver.com/pj3503/222485497313\n오후 3:34\n\n주민등록 주소지 이전변경\n주민등록 주소지 이전 변경 방법 인터넷으로 할수있습니다.\nblog.naver.com\n\n국민취업지원...신청 방법.pdf\n유효기간 ~2023. 5. 11. / 용량 387 KB\n\n주소지 먼저 옮기고 국민취업지원제도 1유형 선발형(청년)으로 신청해주세요~\n오후 3:34''',12.5)
d2=new_after(d,'CHAPTER 3. 상담 스피치','카카오톡 활용 · 원문 3','23-3');exact_page(d2,'''https://blog.naver.com/pj3503/222485497313\n주민등록 주소지 이전변경\n주민등록 주소지 이전 변경 방법 인터넷으로 할수있습니다.\nblog.naver.com\n\n국민취업지원...신청 방법.pdf\n유효기간 ~2023. 5. 29. / 용량 387 KB\n\n내용 확인하고 신청하다가 헷갈리시면 말씀해주세요!\n오후 6:57\n\n★김도형(기1개)\n네 해볼게요\n오후 7:07\n\n★김도형(기1개)\n신청은 해놨어요\n오후 7:29''',13)

# s24 exact chat and payment records.
sl=find('기존 수강생 수급내역');rebuild(sl,'CHAPTER 3. 상담 스피치','기존 수강생 수급내역 · 원문','24')
exact_page(sl,'''★김효령(마야)\n멘토님 이거 들어온 것 맞나요??\n오후 5:53\n\n★김효령(마야)\n3월 수당이욤..!\n오후 5:54''',14)
native_table(sl,92,292,775,170,['구분','일자·시간','입금액','잔액'],[
 ['광주고용청','[표시 없음]','500,000원','789,353원'],['국민취업지원','3월 22일 16:15','500,000원','585,267원'],['훈련비고용부','5월 17일 16:59','200,000원','380,176원']],12,[190,200,180,205])

# s25 exact consultation records.
sl=find('국취제 과정 확대');rebuild(sl,'CHAPTER 3. 상담 스피치','국취제 과정 확대 · 원문 1','25-1')
exact_page(sl,'''이름 : 이수진\n휴대폰 : 010-5095-1761\n나이 : 20대\n성별 : ○남 ◉여\n전공유무 : 비전공자\n공부목적 : □ 진학 ■ 취업 □ 창업 □ 취미\n문의과목 : 구직자, 전산세무 회계 / 오피스 자격증\n희망교육장 : 대전\n당직여부 : ■ 당직\n\n컨택결과 : 재컨택요망 02.28\n*컨택자 : 김호영\n*현재 상황 : 퇴업자\n*공부 목적+니즈(구체적으로) : 모스 자격증 취득\n*전공 유무 및 해본적 있는지 : 없음\n*거주지or직장위치+학원까지 이동시간 : 미파악\n*평일/주말 스케줄이 어떤지 (가용시간 체크) : 상관 없음\n*공부 시작 시기 : 3월\n*생각 중인 공부 기간 : 길지 않음\n*과정 안내 내용 : 모스과정 안내\n*상담 명분 제시+상담 유도 : 국비조회\n*같이 방문하실 지인 : 친한 동생''',11.5)
d=new_after(sl,'CHAPTER 3. 상담 스피치','국취제 과정 확대 · 원문 2','25-2');exact_page(d,'''출처경로 : www.sbsart.com 국비지원과정소개\n출처담당 : 김홍구5 (1588-5530)\n상담신청일 : 2023-02-27 오후 6:37:02\n담당매니저 : 김호영2\n학과구분 : 자격증\n\n컨택내용\n이제 막 졸업한 분으로 친한 동생이랑 모스를 배우고자 함 -> 동생은 4학년\n국비지원이 되는지 물어보았고 둘 다 가능하다고 안내\n모스 자격증을 취득하려하고 특별하게 취득해야할 이유는 없다고 함\n따라서 취업지원까지 가능한 국비제도 안내해드리기로 하였음\n본인은 시간이 많아서 괜찮은데 동생이랑 시간 맞춰보고 상담일정 잡아보겠다고 함\n**수강료 질문에 대한 응대 : X\n**시간표 질문에 대한 응대 : 오후, 저녁\n*최소 과정 : 모스\n*최대 과정 : 사무패키지\n*특이 사항 : X\n*통화 시간 : 6분''',11.5)
d2=new_after(d,'CHAPTER 3. 상담 스피치','국취제 과정 확대 · 원문 3','25-3');exact_page(d2,'''컨택결과 : 완납 3,330,000\n*상담자 : 김호영\n*직업 : 구직자\n*나이 : 00년생\n*거주지 : 도안동\n*과정 관심 계기 : 창업전 취업\n*관심 및 희망 과정 : 모스\n*수강 목적 : 취업 시 도움이 될까\n*학교&직장 내 특이사항 : X\n*평소 가용 가능한 시간 : 저녁 선호 (황예림수강생이랑 같은 시간대 듣기 위해)\n*희망 개강 시기 : 3월\n*희망 타임 : 저녁\n*수강 예산 : 333\n*성격 및 성향 : 보통\n*결제권자/완납 예정 : 아버님 / 완납\n*환불 리스크 : X\n*가족 등록 가망 과목/시기 : X\n*추가 상담 내용 : X''',11.8)
d3=new_after(d2,'CHAPTER 3. 상담 스피치','국취제 과정 확대 · 원문 4','25-4');exact_page(d3,'''이외 상담 내용 :\n건설안전용품을 판매하는 사업을 하고 싶음\n아버님께서는 건설철거 관련해서 사업을 하시는 듯 -> 학생 때 국장 못받음\n국취제 받기 위해 주소이전 안내하였고 서울에 있는 언니집으로 옮기는 쪽으로 할 것 같음\n건축, 경리, 마케팅 분야를 안내하다가 -> 마케팅쪽으로 유도하여 웹디자인 안내\n황예림 수강생이 친구인데 저녁이랑 주말로 수업 다니고 있어서 같은 타임으로 공부하길 희망함\n웹디자인으로 완납\n*이번 상담 때 나의 반성 및 고칠 점 : X''',13)

# s26 exact three records.
sl=find('대상자가 아닐 때');rebuild(sl,'CHAPTER 3. 상담 스피치','대상자가 아닐 때 · 원문 1','26-1')
exact_page(sl,'''컨택결과 : 상담예정 09.20 ★13:30\n컨택일시 : 2022-09-19 12:12:16\n컨택자 : 2-1팀 김호영2\n*컨택자 : 김호영\n*본인확인 및 문의여부 체크 : 본인\n*학생분 현재 상황 질문 : 휴학생\n*배우고 싶은 시기 : 11월\n\n*이외 상담내용\n24세로 건설관련 학과였는데 현장직이 아닌 사무실내에서 근무하는 것을 해보고 싶어 현재 1년간 휴학중\n아직까지 복학할 계획은 없다고 함\n웹디자인과 광고디자인 둘 다 관심이 있는 상황이고 국비지원은 휴학생 불가라고 안내 -> 교육청지원 안내\n본인도 국비지원이 더 많이 되는거냐고 물어봐서 그렇다고 하였음\n부모님과도 조율 중이라 11월 희망한다고 하였는데 상담왔을 시 굳혀버릴 것''',12)
d=new_after(sl,'CHAPTER 3. 상담 스피치','대상자가 아닐 때 · 원문 2','26-2');exact_page(d,'''휴학생으로 부사관학과인데 그쪽으로 취업하지 않을 거라 크게 관심 없다고 함\n친구가 웹프로그래머로 근무하는데 괜찮아 보여서 관심이 생겼고 웹디자인 전망도 좋은 부분 안내\n마지막엔 인테리어쪽도 조금 궁금해하여 고민이 다소 있는 듯함\n익일까지 고민하기로 했고 국취제 참여하는 쪽으로 고려중\n자부담금 설명해줬더니 생각보다 금액이 많이 안나왔다고 하였음\n\n학교 자퇴하지 않고 다니기로 하였고\n따라서 국취제는 내년 9월에 신청하기로 함\n1개월에 1개씩 교육과정 진행하기로 하였음\n김호영2 (7696122)\n기타\n- COD 기한 : 09.21 (수)\n- COD 결제수단 : 카드예정\n- 결제권자 : 부모님\n- 사유 : 국취제 참여 위해 학교 자퇴 고려중''',12.5)

# s27 exact consultation and tuition tables.
sl=find('국취제 할인율');rebuild(sl,'CHAPTER 3. 상담 스피치','국취제 할인율 · 원문 상담','27-1')
exact_page(sl,'''컨택결과 : 완납 3,000,000\n컨택일시 : 2023-05-15 20:04:25\n컨택자 : 2-1팀 김호영2\n*상담자 : 김호영\n*직업 : 구직자\n*나이 : 26세\n*거주지 : 중구 유천동\n*과정 관심 계기 : 취업\n*관심 및 희망 과정 : 컴활, 기계\n*수강 목적 : 취업\n*학교&직장 내 특이사항 : X\n*평소 가용 가능한 시간 : 프리함\n*희망 개강 시기 : 6월\n*희망 타임 : 오후 (9시부터 아니면 됨)\n*희망 예산 : 300\n*성격 및 성향 : 보통\n*결제권자/완납 예정 : 부모님 / 완납\n*환불 리스크 : X\n*추가 등록 가망 과목/시기 : X''',11.5)
d=new_after(sl,'CHAPTER 3. 상담 스피치','국취제 할인율 · 원문 상담내용','27-2');exact_page(d,'''*이외 상담 내용 :\n어머님과 함께 방문하였음\n아들이 졸업하고 취업준비를 안하고 게임만해서 답답해서 전화주심\n아들은 일단 컴활 먼저 하고 나중에 봐서 하겠다는 식이었으나 지금 해야한다고 압박을 주었음\n어머님도 멘토편이었고 조심스럽게 달래가면서 이야기하였음\n결국은 본인도 기계학과를 살려야겠다며 기사 자격증을 준비하는 쪽으로 이야기함\n컴활부터 진행해서 국취제를 통해 기계 분야 공부 진행해보기로 하였음''',13.5)
d2=new_after(d,'CHAPTER 3. 상담 스피치','국취제 할인율 · 원문 수강료','27-3');native_table(d2,42,105,875,190,['NO','차수','과목순번','과목','수강료'],[['1','1','7087','포토샵','350,000'],['2','1','7103','캐드1','350,000'],['3','1','7102','컴활','300,000'],['4','1','8123','캐드2','350,000'],['5','1','7252','라이노1(방학)','350,000']],11.5,[75,75,160,350,215]);box(d2,42,315,875,30,'주말반',PALE2,BORDER,13,True,DARK,PP_ALIGN.CENTER);native_table(d2,42,355,875,130,['NO','차수','과목순번','과목','수강료'],[['1','1','7534','인벤터','400,000'],['2','1','16764','라이노2/주말','400,000'],['3','1','21793','일반기계기사/주말','450,000']],11.5,[75,75,160,350,215])

# Remove all remaining non-cover pictures and enforce font.
for si,sl in enumerate(prs.slides,1):
 for s in list(sl.shapes):
  if si>1 and s.shape_type==13:rm(s)
  if getattr(s,'has_text_frame',False):
   for p in s.text_frame.paragraphs:
    for r in p.runs:r.font.name=FONT
  if getattr(s,'has_table',False):
   for row in s.table.rows:
    for c in row.cells:
     for p in c.text_frame.paragraphs:
      for r in p.runs:r.font.name=FONT

prs.save(out);print(f'saved={out} slides={len(prs.slides)}')
