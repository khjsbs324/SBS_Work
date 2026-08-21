import sys
from pathlib import Path
sys.path.insert(0,str(Path(__file__).resolve().parents[1]/'pptdeps'))
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN,MSO_ANCHOR
from pptx.dml.color import RGBColor

src,out=map(Path,sys.argv[1:3]);prs=Presentation(src)
FONT='Noto Sans KR';NAVY=RGBColor(6,70,130);BLUE=RGBColor(10,79,153);DARK=RGBColor(10,53,106);PALE=RGBColor(238,247,255);PALE2=RGBColor(221,240,255);BORDER=RGBColor(185,217,245);TEXT=RGBColor(30,43,56);WHITE=RGBColor(255,255,255);WARM=RGBColor(255,245,223)

def title_of(sl):
    candidates=[]
    for sh in sl.shapes:
        if getattr(sh,'has_text_frame',False) and sh.text.strip() and sh.top.pt<90:
            if sh.name=='Title' or 'Title' in sh.name:candidates.append(sh.text.strip())
    return max(candidates,key=len) if candidates else ''
def all_text(sl):
    parts=[]
    for sh in sl.shapes:
        if getattr(sh,'has_text_frame',False):parts.append(sh.text)
        if getattr(sh,'has_table',False):
            for row in sh.table.rows:
                parts.extend(c.text for c in row.cells)
    return ''.join(parts).strip()
def delete_slide(idx):
    slide_id=prs.slides._sldIdLst[idx];rel=prs.part.rels[slide_id.rId];prs.part.drop_rel(slide_id.rId);del prs.slides._sldIdLst[idx]

# Remove all duplicated source-image pages and the original 28-50 image reference block.
for i in range(len(prs.slides)-1,0,-1):
    sl=prs.slides[i];t=title_of(sl)
    if '원본' in t or not all_text(sl):delete_slide(i)

# Restore clear spacing between the small CONTENTS label and the slide-2 title.
for sh in prs.slides[1].shapes:
    if getattr(sh,'has_text_frame',False) and '국비 컨택' in sh.text and '&' in sh.text:
        sh.top=Pt(39)

def rm(sh):sh._element.getparent().remove(sh._element)
def clear_body(sl):
    for sh in list(sl.shapes):
        if sh.top.pt>82:rm(sh)
def text(sl,x,y,w,h,txt,size=14,bold=False,color=TEXT,align=PP_ALIGN.LEFT):
    sh=sl.shapes.add_textbox(Pt(x),Pt(y),Pt(w),Pt(h));tf=sh.text_frame;tf.clear();tf.word_wrap=True;tf.margin_left=tf.margin_right=Pt(2);tf.margin_top=tf.margin_bottom=Pt(1);tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    for i,line in enumerate(str(txt).split('\n')):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph();p.alignment=align;r=p.add_run();r.text=line;r.font.name=FONT;r.font.size=Pt(size);r.font.bold=bold;r.font.color.rgb=color
    return sh
def box(sl,x,y,w,h,txt='',fill=WHITE,line=BORDER,size=14,bold=False,color=TEXT,align=PP_ALIGN.LEFT):
    sh=sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Pt(x),Pt(y),Pt(w),Pt(h));sh.fill.solid();sh.fill.fore_color.rgb=fill;sh.line.color.rgb=line;sh.line.width=Pt(1)
    tf=sh.text_frame;tf.clear();tf.word_wrap=True;tf.margin_left=tf.margin_right=Pt(14);tf.margin_top=tf.margin_bottom=Pt(6);tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    for i,line in enumerate(str(txt).split('\n')):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph();p.alignment=align;r=p.add_run();r.text=line;r.font.name=FONT;r.font.size=Pt(size);r.font.bold=bold;r.font.color.rgb=color
    return sh
def table(sl,x,y,w,h,headers,rows,size=12,widths=None):
    shp=sl.shapes.add_table(len(rows)+1,len(headers),Pt(x),Pt(y),Pt(w),Pt(h));tb=shp.table
    if widths:
        for i,v in enumerate(widths):tb.columns[i].width=Pt(v)
    for ri,row in enumerate([headers]+rows):
        for ci,v in enumerate(row):
            c=tb.cell(ri,ci);c.fill.solid();c.fill.fore_color.rgb=NAVY if ri==0 else (PALE if ri%2 else WHITE);c.margin_left=c.margin_right=Pt(5);c.margin_top=c.margin_bottom=Pt(4);tf=c.text_frame;tf.clear();tf.word_wrap=True;tf.vertical_anchor=MSO_ANCHOR.MIDDLE
            p=tf.paragraphs[0];p.alignment=PP_ALIGN.CENTER;r=p.add_run();r.text=str(v);r.font.name=FONT;r.font.size=Pt(size);r.font.bold=(ri==0 or ci==0);r.font.color.rgb=WHITE if ri==0 else TEXT
    return shp
def find_title(title):return next(sl for sl in prs.slides if title_of(sl)==title)

# Convert former evidence-image topics into native, editable structures.
sl=find_title('카카오톡 활용');clear_body(sl)
box(sl,42,96,875,38,'상담 시 카톡 활용',PALE2,BORDER,15,True,DARK,PP_ALIGN.CENTER)
table(sl,62,151,835,292,['단계','활용 내용','상담 포인트'],[
 ['상담 전','신청 링크와 안내 자료 전달','문의자가 필요한 정보를 한 화면에서 확인'],
 ['조건 확인','주소지 이전·정부24 등 절차 안내','질문과 답변을 대화 기록으로 남김'],
 ['서류 진행','PDF와 신청 서류 전달','제출 여부와 누락 자료를 확인'],
 ['상담 후','진행 상태와 추가 질문 관리','상담 흐름이 끊기지 않도록 후속 관리']],12,[140,330,365])
box(sl,142,463,675,32,'링크·자료·답변을 한 대화 흐름으로 관리',WARM,BORDER,12.5,True,DARK,PP_ALIGN.CENTER)

sl=find_title('기존 수강생 수급내역');clear_body(sl)
box(sl,42,96,875,38,'기존 수강생 수급내역',PALE2,BORDER,15,True,DARK,PP_ALIGN.CENTER)
table(sl,92,160,775,230,['구분','확인 일자','입금액','화면상 잔액'],[
 ['광주고용지원','3월 22일','500,000원','789,353원'],
 ['광주고용지원','확인 화면','500,000원','585,267원'],
 ['훈련비고용부','5월 17일','200,000원','380,167원']],13,[190,170,190,225])
box(sl,142,423,675,55,'실제 입금 사례를 근거로 지원금 수령 흐름을 설명하되, 개인별 지급 조건과 금액은 달라질 수 있음을 함께 안내합니다.',PALE,BORDER,12.5,False,TEXT,PP_ALIGN.CENTER)

sl=find_title('국취제 과정 확대');clear_body(sl)
box(sl,42,105,410,130,'현재 문의\n단순 컴활 취업',WARM,BORDER,18,True,DARK,PP_ALIGN.CENTER)
box(sl,508,105,410,130,'확장 방향\n정규과정으로 연결',PALE2,BORDER,18,True,DARK,PP_ALIGN.CENTER)
text(sl,456,148,48,45,'→',28,True,BLUE,PP_ALIGN.CENTER)
table(sl,82,280,795,150,['질문','구체화 방향','추천 과정'],[
 ['사무직을 원할 때','회계·정산 중심','경리'],
 ['홍보·콘텐츠에 관심 있을 때','채널 운영·광고 중심','마케팅']],13,[245,275,275])
box(sl,142,455,675,32,'넓은 취업 목표를 실제 직무와 정규 커리큘럼으로 구체화',PALE,BORDER,12.5,True,DARK,PP_ALIGN.CENTER)

sl=find_title('대상자가 아닐 때');clear_body(sl)
table(sl,62,112,835,270,['상황','가능한 안내','금지 사항'],[
 ['휴학생으로 국비 대상이 아님','교육청 지원(할인) 등 다른 지원 방식 확인','국비 참여를 위해 자퇴를 권하지 않음'],
 ['현재 조건으로 지원이 어려움','일반과정·장학지원·일정 조정 검토','가능 여부를 단정하거나 과장하지 않음'],
 ['추후 조건 변경 가능','재학·취업 상태 변경 시 재확인','현재 등록을 무리하게 압박하지 않음']],13,[250,335,250])
box(sl,112,417,735,62,'핵심 원칙\n대상자가 아니어도 학습 목표에 맞는 합법적이고 현실적인 대안을 안내합니다.',WARM,BORDER,14,True,DARK,PP_ALIGN.CENTER)

sl=find_title('국취제 할인율');clear_body(sl)
box(sl,82,110,350,110,'잘못된 접근\n국취제 = 할인',WARM,BORDER,18,True,DARK,PP_ALIGN.CENTER)
box(sl,528,110,350,110,'올바른 접근\n국취제 = 취업지원 제도',PALE2,BORDER,18,True,DARK,PP_ALIGN.CENTER)
text(sl,448,145,64,45,'→',28,True,BLUE,PP_ALIGN.CENTER)
table(sl,82,268,795,154,['안내 기준','설명'],[
 ['비용','할인 여부만이 아니라 실제 자부담과 지급 구조를 설명'],
 ['과정','취업 목표와 과정 적합성을 우선 확인'],
 ['상담','국비지원은 할인이 필수가 아니라는 점을 명확히 안내']],13,[210,585])
box(sl,142,456,675,32,'제도의 목적과 과정 적합성을 중심으로 상담',PALE,BORDER,12.5,True,DARK,PP_ALIGN.CENTER)

# Enforce Noto Sans KR on every editable text object.
for sl in prs.slides:
    for sh in sl.shapes:
        if getattr(sh,'has_text_frame',False):
            for p in sh.text_frame.paragraphs:
                for r in p.runs:r.font.name=FONT
        if getattr(sh,'has_table',False):
            for row in sh.table.rows:
                for c in row.cells:
                    for p in c.text_frame.paragraphs:
                        for r in p.runs:r.font.name=FONT

prs.save(out);print(f'saved={out} slides={len(prs.slides)}')
