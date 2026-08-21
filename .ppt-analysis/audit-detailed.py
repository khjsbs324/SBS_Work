import sys
sys.stdout.reconfigure(encoding='utf-8')
from pathlib import Path
sys.path.insert(0,str(Path(__file__).resolve().parents[1]/'pptdeps'))
from pptx import Presentation

p=Path(sys.argv[1]);prs=Presentation(p);sw,sh=prs.slide_width,prs.slide_height
issues=[];small=[];titles=[]
for si,sl in enumerate(prs.slides,1):
    title=''
    pics=0;tables=0
    for sp in sl.shapes:
        if sp.name=='Title' or ('Title' in sp.name and sp.top < sh*.2):title=getattr(sp,'text','').replace('\n',' ')
        if sp.shape_type==13:pics+=1
        if getattr(sp,'has_table',False):tables+=1
        if sp.left < -1000 or sp.top < -1000 or sp.left+sp.width > sw+1000 or sp.top+sp.height > sh+1000:issues.append((si,sp.name,'out_of_bounds'))
        if getattr(sp,'has_text_frame',False):
            for para in sp.text_frame.paragraphs:
                for r in para.runs:
                    if r.font.size and r.font.size.pt < 11 and sp.top > sh*.15 and 'Footer' not in sp.name:small.append((si,sp.name,round(r.font.size.pt,1),r.text[:30]))
        if getattr(sp,'has_table',False):
            for row in sp.table.rows:
                for cell in row.cells:
                    for para in cell.text_frame.paragraphs:
                        for r in para.runs:
                            if r.font.size and r.font.size.pt < 10.5:small.append((si,sp.name,round(r.font.size.pt,1),r.text[:30]))
    titles.append((si,title,pics,tables,len(sl.shapes)))
print('slides',len(prs.slides),'out_of_bounds',len(issues),'small_runs',len(small))
for x in issues[:30]:print('BOUND',x)
for x in small[:50]:print('SMALL',x)
for x in titles:print('SLIDE',*x,sep='\t')
