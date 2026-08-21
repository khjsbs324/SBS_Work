import sys
from pathlib import Path
sys.path.insert(0,str(Path(__file__).resolve().parents[1]/'pptdeps'))
from pptx import Presentation

p=Path(sys.argv[1]);prs=Presentation(p);sw,sh=prs.slide_width,prs.slide_height
small=[];bounds=[];overlaps=[];pictures=[]
for si,sl in enumerate(prs.slides,1):
    text_shapes=[]
    for sp in sl.shapes:
        if sp.shape_type==13:pictures.append((si,sp.name))
        if sp.left<0 or sp.top<0 or sp.left+sp.width>sw or sp.top+sp.height>sh:bounds.append((si,sp.name))
        has_text=getattr(sp,'has_text_frame',False) and sp.text.strip()
        if has_text:text_shapes.append(sp)
        if getattr(sp,'has_text_frame',False):
            for para in sp.text_frame.paragraphs:
                for r in para.runs:
                    if r.font.size and r.font.size.pt<11 and sp.top.pt>82 and 'Footer' not in sp.name:small.append((si,sp.name,round(r.font.size.pt,1),r.text[:30]))
        if getattr(sp,'has_table',False):
            for row in sp.table.rows:
                for cell in row.cells:
                    for para in cell.text_frame.paragraphs:
                        for r in para.runs:
                            if r.font.size and r.font.size.pt<10.5:small.append((si,sp.name,round(r.font.size.pt,1),r.text[:30]))
    for i,a in enumerate(text_shapes):
        ax1,ay1,ax2,ay2=a.left,a.top,a.left+a.width,a.top+a.height
        for b in text_shapes[i+1:]:
            bx1,by1,bx2,by2=b.left,b.top,b.left+b.width,b.top+b.height
            ix=max(0,min(ax2,bx2)-max(ax1,bx1));iy=max(0,min(ay2,by2)-max(ay1,by1))
            if ix*iy>0:
                # Ignore deliberate containment (e.g. text-bearing badge within a large card is not used in v4 tables).
                a_contains=ax1<=bx1 and ay1<=by1 and ax2>=bx2 and ay2>=by2
                b_contains=bx1<=ax1 and by1<=ay1 and bx2>=ax2 and by2>=ay2
                if not(a_contains or b_contains):overlaps.append((si,a.name,b.name))
print(f'slides={len(prs.slides)} pictures={len(pictures)} bounds={len(bounds)} small={len(small)} overlaps={len(overlaps)}')
for x in pictures:print('PICTURE',x)
for x in bounds:print('BOUND',x)
for x in small:print('SMALL',x)
for x in overlaps:print('OVERLAP',x)
